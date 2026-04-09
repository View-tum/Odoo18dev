from __future__ import annotations

from pathlib import Path
from textwrap import dedent
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(r"C:\365_project\TheCool18e\Dev\reports\training_manufacturing_20260331")

PPTX_OUT = BASE / "manufacturing_uat_replenishment_route_20260402.pptx"
SLIDE_MD_OUT = BASE / "08_uat_replenishment_route_slides_th.md"
FLOW_MD_OUT = BASE / "09_fg_pss_th_01005_flow_diagram_th.md"
FLOW_SVG_OUT = BASE / "09_fg_pss_th_01005_flow_diagram_th.svg"
TRAINER_OUT = BASE / "10_uat_replenishment_route_trainer_script_th.md"


SLIDES = [
    {
        "title": "Replenishment, Route, Rule, Operation Type และ Putaway ใน UAT",
        "subtitle": "อธิบายจาก config และ master data จริงของฐาน uat",
        "bullets": [
            "ใช้ตัวอย่างจริง: FG-PNC-TH-01001, FG-PSS-TH-01005, SM-PLS-UP-01001, SO-PSS-LO-01001",
            "ชี้ความต่างระหว่างตัวจุดชนวน, policy, rule ย่อย, เอกสารจริง และกฎจัดเก็บ",
            "ใช้เป็นสไลด์ประกอบ training หรือเป็น appendix เพิ่มใน Manu training หลักได้",
        ],
    },
    {
        "title": "ภาพรวมคำ 5 คำในระบบ",
        "bullets": [
            "Replenishment = ตัวเริ่ม demand เช่น orderpoint, shortage, หรือ MTO demand",
            "Route = policy ว่าจะ Buy, Manufacture, Transfer หรือ MTO",
            "Rule = คำสั่งย่อยจริง ระบุ source, destination, action และ operation type",
            "Operation Type = ชนิดเอกสารที่ผู้ใช้เห็นจริง เช่น Receipts, INT-PL, INT-PH, MO-PH, MO-PL",
            "Putaway = กฎจัดเก็บหลัง stock ถึง location ใหญ่แล้ว",
        ],
    },
    {
        "title": "Replenishment ใน uat ทำงานอย่างไร",
        "bullets": [
            "FG-PNC-TH-01001 มี orderpoint id 175 ที่ GMP/Stock, min 0, max 0, trigger auto",
            "FG-PSS-TH-01005 มี orderpoint id 249 ที่ GMP/Stock, route Manufacture, trigger auto",
            "FG-PSS-TH-01005 มี orderpoint id 308 ที่ M-WH/Stock ด้วย จึงมี replenishment ได้มากกว่าหนึ่งคลัง",
            "min/max แบบ 0/0 ในฐานนี้ใช้แนวคิดเติม shortage กลับขึ้นมาที่ 0 ไม่ได้แปลว่าไม่เติมเลย",
        ],
    },
    {
        "title": "Route จริงที่ใช้ใน uat",
        "bullets": [
            "Buy (id 5) ใช้กับของที่ต้องซื้อเข้า เช่น packaging และ RM หลายตัว",
            "Manufacture (id 6) เป็น route มาตรฐานที่ orderpoint บางตัวเรียกใช้",
            "Manufacture (Pharma) (id 62) และ Manufacture (Plastic) (id 63) ใช้แยกโรงงานและเลขเอกสาร MO",
            "Auto Transfer Semi (Pharma) (id 60) และ Auto Transfer Semi (Plastic) (id 61) ใช้ดึง semi ก่อนหรือระหว่าง chain",
            "Replenish on Order (MTO) (id 1) เป็น route ที่รวมหลาย rule ตามปลายทางไว้ใน route เดียว",
        ],
    },
    {
        "title": "Rule คือจุดที่ logic ลงรายละเอียดจริง",
        "bullets": [
            "Rule 7: GMP: Stock (Buy) -> action buy, destination GMP/Stock, operation type Receipts",
            "Rule 146: GMP: Stock (Production) (copy) (Pharma) -> action manufacture, source GMP/Stock/คลังลอย, destination GMP/Stock, operation type Manufacturing Pharma",
            "Rule 147: GMP: Stock (Production) (copy) (Plastic) -> action manufacture, source GMP/Stock/คลังลอย, destination GMP/Stock, operation type Manufacturing Plastic",
            "Rule 144 และ 145 เป็น pull rules สำหรับ Transfer Pharma / Transfer Plastic เพื่อดึง semi เข้าคลังลอย",
        ],
    },
    {
        "title": "Operation Type คือเอกสารจริงที่คนหน้างานเห็น",
        "bullets": [
            "Receipts -> เอกสารรับของจาก vendor",
            "Transfer Pharma / Transfer Plastic -> เอกสาร internal transfer คนละฝั่งโรงงาน",
            "Manufacturing Pharma / Manufacturing Plastic -> MO คนละ sequence และคนละ operation type",
            "ใน training ต้องให้ทีมดู Operation Type เป็นหลัก ไม่ใช่ดูจากชื่อสินค้าอย่างเดียว",
        ],
    },
    {
        "title": "Putaway ใน uat รับไม้ต่อหลัง document เสร็จ",
        "bullets": [
            "ถ้าเข้า GMP/Stock และ category เป็น RM/พลาสติก -> ไป GMP/Stock/RM/PL01",
            "ถ้าเข้า GMP/Stock และ category เป็น RM/สารเคมี -> ไป GMP/Stock/RM/สารเคมี",
            "ถ้าเข้า GMP/Stock และ category เป็น SFG/พลาสติก -> ไป GMP/Stock/Semi/พลาสติก",
            "ถ้าเข้า GMP/Stock และ category เป็น SFG/น้ำยา -> ไป GMP/Stock/Semi/โรงงานยา",
            "Putaway ไม่ได้สร้าง PO หรือ MO แต่ใช้จัดเก็บต่อหลัง stock move สำเร็จ",
        ],
    },
    {
        "title": "ตัวอย่าง FG-PNC-TH-01001",
        "bullets": [
            "orderpoint id 175 เป็นตัว trigger ที่ GMP/Stock",
            "route บน product เป็น Manufacture (Pharma)",
            "orderpoint ไม่ล็อก route เอง จึงปล่อยให้ product route เป็นคนบอกว่าจะไปโรงงานยา",
            "เมื่อ demand เกิด ระบบจึงสร้าง chain ฝั่ง Manufacturing Pharma",
        ],
    },
    {
        "title": "ตัวอย่าง FG-PSS-TH-01005",
        "bullets": [
            "trigger มาจาก orderpoint 249 ที่ GMP/Stock",
            "orderpoint route = Manufacture แต่ product route = Manufacture (Pharma)",
            "ชั้นบนสุดของ BoM มี FG-PSS-TH-02001 x16, PK-CAR-PS-01003 x1 และเทปกาว x1",
            "เมื่อแตกต่อ จะมีทั้ง branch ที่ Buy, branch ที่เป็น Pharma, และ branch ที่เป็น Plastic อยู่ใน chain เดียวกัน",
        ],
    },
    {
        "title": "ตัวอย่างชั้นลึกของ FG-PSS-TH-01005",
        "bullets": [
            "FG-PSS-TH-02001 แตกเป็น FG-PSS-TH-03001 x160, PK-BOX-PS-01002 และ PK-SHF-PS-01003",
            "FG-PSS-TH-03001 แตกเป็น FG-PSS-TH-04001 ถึง FG-PSS-TH-04006 และวัสดุบางตัวที่ Buy",
            "ตัวอย่าง FG-PSS-TH-04001 มีส่วนประกอบทั้ง SM-PSS-TH-02001, SM-PSS-TH-02002, SM-JOI-PK-02002, RM-FIL-PS-01004 และ SO-PSS-LO-01001",
            "จุดนี้เองที่ทำให้ chain เดียวไหลทั้ง Plastic, Pharma และ Buy",
        ],
    },
    {
        "title": "สรุป logic ที่ควรใช้สอนทีม",
        "bullets": [
            "Replenishment เป็นตัวเริ่มเรื่อง แต่ไม่ได้ตัดสินทุกอย่างด้วยตัวเอง",
            "Route เป็น policy, Rule เป็น logic ย่อยจริง, Operation Type เป็นเอกสารจริง",
            "Putaway ทำงานตอนท้ายเพื่อจัดเก็บ ไม่ได้เป็น procurement engine",
            "เวลาระบบไม่ไปต่อ ให้เช็กตามลำดับ: trigger -> route -> rule -> operation type -> putaway",
        ],
    },
]


SLIDE_MD = dedent(
    """
    # Slide ภาษาไทย: Replenishment, Route, Rule, Operation Type และ Putaway ใน UAT

    ## สไลด์ 1: หัวข้อ
    - Replenishment, Route, Rule, Operation Type และ Putaway ใน UAT
    - อธิบายจาก config และ master data จริงของฐาน uat
    - ใช้เป็น appendix เพิ่มในชุด training manufacturing หลัก

    ## สไลด์ 2: ภาพรวมคำ 5 คำในระบบ
    - Replenishment = ตัวเริ่ม demand
    - Route = policy ว่าจะ Buy, Manufacture, Transfer หรือ MTO
    - Rule = คำสั่งย่อยจริงที่ระบุ source, destination, action และ operation type
    - Operation Type = เอกสารจริงที่ผู้ใช้เห็น
    - Putaway = กฎจัดเก็บหลัง stock ถึงปลายทาง

    ## สไลด์ 3: Replenishment ใน uat ทำงานอย่างไร
    - FG-PNC-TH-01001 มี orderpoint id 175 ที่ GMP/Stock, trigger auto, min 0, max 0
    - FG-PSS-TH-01005 มี orderpoint id 249 ที่ GMP/Stock และ id 308 ที่ M-WH/Stock
    - min/max แบบ 0/0 ในฐานนี้ใช้เติม shortage กลับขึ้นมาที่ 0
    - Replenishment จึงเป็น trigger ของ procurement ไม่ใช่เอกสารปลายทาง

    ## สไลด์ 4: Route จริงที่ใช้ใน uat
    - Buy (id 5)
    - Manufacture (id 6)
    - Manufacture (Pharma) (id 62)
    - Manufacture (Plastic) (id 63)
    - Auto Transfer Semi (Pharma) (id 60)
    - Auto Transfer Semi (Plastic) (id 61)
    - Replenish on Order (MTO) (id 1)

    ## สไลด์ 5: Rule คือจุดที่ logic ลงรายละเอียดจริง
    - Rule 7: GMP: Stock (Buy)
    - Rule 146: GMP: Stock (Production) (copy) (Pharma)
    - Rule 147: GMP: Stock (Production) (copy) (Plastic)
    - Rule 144 และ 145: pull rules สำหรับ Transfer Pharma / Transfer Plastic

    ## สไลด์ 6: Operation Type คือเอกสารจริง
    - Receipts
    - Transfer Pharma
    - Transfer Plastic
    - Manufacturing Pharma
    - Manufacturing Plastic
    - ต้องให้ทีมดู Operation Type เพื่อแยกฝั่งงานให้ถูก

    ## สไลด์ 7: Putaway ใน uat
    - RM/พลาสติก -> GMP/Stock/RM/PL01
    - RM/สารเคมี -> GMP/Stock/RM/สารเคมี
    - SFG/พลาสติก -> GMP/Stock/Semi/พลาสติก
    - SFG/น้ำยา -> GMP/Stock/Semi/โรงงานยา
    - Putaway ทำงานหลัง document เสร็จแล้ว

    ## สไลด์ 8: ตัวอย่าง FG-PNC-TH-01001
    - orderpoint 175 เป็น trigger
    - product route เป็น Manufacture (Pharma)
    - demand ถูกจุดโดย orderpoint แต่ chain ฝั่งผลิตอิง product route

    ## สไลด์ 9: ตัวอย่าง FG-PSS-TH-01005
    - orderpoint 249 route = Manufacture
    - product route = Manufacture (Pharma)
    - top BOM มี FG-PSS-TH-02001 x16, PK-CAR-PS-01003 x1, เทปกาว x1
    - chain นี้มีทั้ง Pharma, Plastic และ Buy

    ## สไลด์ 10: ชั้นลึกของ FG-PSS-TH-01005
    - FG-PSS-TH-02001 -> FG-PSS-TH-03001 x160 + packaging buy
    - FG-PSS-TH-03001 -> FG-PSS-TH-04001..04006 + buy items
    - สีแต่ละตัวแตกลง semi plastic และ solution pharma ต่อ

    ## สไลด์ 11: สรุป logic ที่ควรใช้สอนทีม
    - Replenishment เป็นตัวเริ่มเรื่อง
    - Route เป็น policy
    - Rule เป็น logic ย่อยจริง
    - Operation Type เป็นเอกสารจริง
    - Putaway เป็น storage logic ตอนท้าย
    """
).strip() + "\n"


TRAINER_MD = dedent(
    """
    # Trainer Script ภาษาไทย: Replenishment, Route, Rule, Operation Type และ Putaway ใน UAT

    ## วัตถุประสงค์
    ใช้ script นี้เมื่อผู้สอนต้องการอธิบายว่า demand 1 ตัวในระบบ Odoo ของเราไหลจาก trigger ไปจนถึง document และ location ย่อยอย่างไร โดยอิงจาก config จริงในฐาน `uat`

    ---

    ## เปิดเรื่อง
    วันนี้เราจะอธิบายคำ 5 คำที่คนใช้ระบบสับสนกันบ่อย คือ Replenishment, Route, Rule, Operation Type และ Putaway โดยจะใช้ข้อมูลจริงจากฐาน uat ไม่ใช่ตัวอย่างสมมติ เป้าหมายคือให้ทุกคนตอบได้ว่าถ้าของขาด ระบบจะไปผลิต ไปซื้อ หรือไปโอน เพราะอะไร

    ---

    ## สไลด์ 1
    ให้บอกผู้เรียนว่าหัวข้อนี้เป็น appendix สำคัญของ training manufacturing เพราะเวลา flow ไม่ไปต่อ ปัญหามักอยู่ที่ 5 คำนี้ ไม่ใช่ที่หน้าจอปลายทาง

    ## สไลด์ 2
    อธิบายความแตกต่างของ 5 คำให้ชัด
    - Replenishment คือ trigger
    - Route คือ policy
    - Rule คือ logic ย่อยจริง
    - Operation Type คือเอกสารจริง
    - Putaway คือการจัดเก็บหลังงานเสร็จ

    เน้นว่าคำทั้ง 5 ไม่ได้ซ้อนกันแบบคำพ้อง แต่ทำงานคนละชั้น

    ## สไลด์ 3
    ใช้ FG-PNC-TH-01001 และ FG-PSS-TH-01005 เป็นตัวอย่าง
    - FG-PNC-TH-01001 มี orderpoint 175 ที่ GMP/Stock
    - FG-PSS-TH-01005 มี orderpoint 249 ที่ GMP/Stock และ 308 ที่ M-WH/Stock

    จุดที่ต้องอธิบาย:
    - trigger auto หมายถึงระบบยิง procurement ให้เอง
    - min/max 0/0 ในฐานนี้หมายถึงเติม shortage กลับมาที่ 0
    - orderpoint เป็นตัวเริ่ม แต่ยังไม่ใช่ตัวกำหนด document สุดท้ายทั้งหมด

    ## สไลด์ 4
    อธิบาย Route ที่มีอยู่จริงใน UAT
    - Buy
    - Manufacture
    - Manufacture (Pharma)
    - Manufacture (Plastic)
    - Auto Transfer Semi (Pharma)
    - Auto Transfer Semi (Plastic)
    - Replenish on Order (MTO)

    ให้ย้ำว่า Route คือ policy ของสินค้า ไม่ใช่เอกสาร

    ## สไลด์ 5
    อธิบาย Rule ด้วยตัวอย่างจริง
    - Rule 7 ใช้กับ Buy
    - Rule 146 ใช้กับ Manufacturing Pharma
    - Rule 147 ใช้กับ Manufacturing Plastic
    - Rule 144/145 ใช้กับ Transfer Semi

    จุดสำคัญที่ต้องพูด:
    - source location
    - destination location
    - action
    - operation type

    ให้พูดว่าเวลาระบบไปผิดทาง ต้องลงมาดูที่ Rule เป็นหลัก

    ## สไลด์ 6
    ย้ำว่า Operation Type คือสิ่งที่คนใช้งานเห็นเป็น document จริง
    เช่น
    - Receipts
    - Transfer Pharma
    - Transfer Plastic
    - Manufacturing Pharma
    - Manufacturing Plastic

    ให้สอนทีมดู operation type เพื่อแยกโรงงาน ไม่ใช่เดาจากชื่อสินค้าอย่างเดียว

    ## สไลด์ 7
    อธิบาย Putaway ว่าเป็น layer หลังสุด
    ใช้ตัวอย่างจริง:
    - RM/พลาสติก ไป PL01
    - SFG/พลาสติก ไป Semi/พลาสติก
    - SFG/น้ำยา ไป Semi/โรงงานยา

    ประโยคสำคัญ:
    Putaway ไม่ได้สร้าง PO ไม่ได้สร้าง MO แต่เป็นตัวจัดเก็บเมื่อของไปถึง location ใหญ่แล้ว

    ## สไลด์ 8
    อธิบาย FG-PNC-TH-01001
    ลำดับที่ต้องพูด:
    1. orderpoint 175 ยิง demand
    2. product route เป็น Manufacture (Pharma)
    3. ระบบจึงไป chain ฝั่งผลิตยา

    จุดที่ต้องย้ำ:
    orderpoint เป็นตัวเริ่ม แต่ route บนสินค้าเป็นคนกำหนดฝั่งผลิต

    ## สไลด์ 9
    อธิบาย FG-PSS-TH-01005
    - orderpoint 249 route = Manufacture
    - product route = Manufacture (Pharma)
    - top BOM มี FG-PSS-TH-02001 x16, PK-CAR-PS-01003 x1 และเทปกาว x1

    ให้ผู้เรียนเห็นว่า chain นี้มีทั้ง:
    - ผลิต
    - ซื้อ
    - และลึกลงไปยังแยก Plastic กับ Pharma

    ## สไลด์ 10
    ใช้ชั้นลึกของ FG-PSS-TH-01005 เพื่ออธิบายของจริง
    - FG-PSS-TH-02001 แตกเป็น FG-PSS-TH-03001 และ packaging buy
    - FG-PSS-TH-03001 แตกเป็น 6 สี
    - ตัวอย่างสี FG-PSS-TH-04001 ใช้ทั้ง semi plastic, joiner plastic, filter, และ solution pharma

    ตรงนี้ต้องชี้ให้ชัดว่า:
    - Plastic branch ใช้ route 61/63
    - Pharma branch ใช้ route 60/62
    - Buy branch ใช้ route 5

    ## สไลด์ 11
    ปิดด้วยสรุป 5 คำอีกครั้ง
    - Replenishment = ตัวเริ่มเรื่อง
    - Route = บอกทาง
    - Rule = บอกวิธีทำจริง
    - Operation Type = หน้าตาเอกสาร
    - Putaway = บอกว่าพอไปถึงแล้วเก็บตรงไหน

    ประโยคปิด:
    ถ้าระบบไม่สร้างเอกสารที่เราคาด อย่าเริ่มจากโทษหน้าจอปลายทาง ให้ไล่จาก trigger ไป route ไป rule แล้วค่อยดู operation type และ putaway

    ---

    ## คำถามที่คนเรียนมักถาม
    ### ถาม: Replenishment กับ Orderpoint เหมือนกันไหม
    ตอบ: ไม่เหมือนกัน Orderpoint เป็นหนึ่งใน trigger ของ replenishment แต่ replenishment กว้างกว่า เพราะ MTO demand ก็เป็น replenishment ได้

    ### ถาม: ทำไมสินค้าบางตัวมีหลาย route
    ตอบ: เพราะ route คนละตัวรับผิดชอบคนละ behavior เช่น สินค้าหนึ่งตัวอาจต้องมีทั้ง transfer semi และ manufacture

    ### ถาม: ทำไม route Manufacture ที่ orderpoint กับ route Manufacture (Pharma) บนสินค้าไม่เหมือนกัน
    ตอบ: เพราะ orderpoint บอกเชิง procurement ว่าต้องผลิต แต่ route เฉพาะบน product/child chain ใช้แยก logic โรงงานและ operation type

    ### ถาม: Putaway มีผลกับ procurement ไหม
    ตอบ: ไม่มี Putaway มีผลตอนของมาถึงปลายทางแล้วเท่านั้น

    ### ถาม: ถ้าของขาด ระบบจะ Buy หรือ Manufacture ดูตรงไหน
    ตอบ: ดูที่ route และ rule ที่วิ่งจริงของสินค้านั้น รวมทั้ง route บน orderpoint ถ้ามี
    """
).strip() + "\n"


FLOW_MD = dedent(
    """
    # Flow Diagram: FG-PSS-TH-01005 ใน UAT

    ## จุดตั้งต้น
    - Product: `FG-PSS-TH-01005`
    - Product Route: `Manufacture (Pharma)` id `62`
    - Orderpoint: `249`
    - Orderpoint Location: `GMP/Stock`
    - Orderpoint Route: `Manufacture` id `6`
    - Trigger: `auto`

    ## Top-level chain
    1. Orderpoint 249 ตรวจ forecast ที่ `GMP/Stock`
    2. ถ้าของไม่พอ ระบบเริ่ม procurement
    3. Product route บอกว่าปลายทางหลักต้องวิ่งฝั่ง `Manufacture (Pharma)`
    4. Rule `146` ทำให้ document หลักออกเป็น `Manufacturing Pharma`

    ## Top BOM
    - `FG-PSS-TH-02001` x `16`
    - `PK-CAR-PS-01003` x `1` -> `Buy`
    - `เทปกาว` x `1` -> `Buy`

    ## Mid BOM
    `FG-PSS-TH-02001` แตกเป็น
    - `FG-PSS-TH-03001` x `160`
    - `PK-BOX-PS-01002` x `16` -> `Buy`
    - `PK-SHF-PS-01003` x `0.01424` -> `Buy`

    ## Color layer
    `FG-PSS-TH-03001` แตกเป็น
    - `FG-PSS-TH-04001` x `160`
    - `FG-PSS-TH-04002` x `160`
    - `FG-PSS-TH-04003` x `160`
    - `FG-PSS-TH-04004` x `160`
    - `FG-PSS-TH-04005` x `160`
    - `FG-PSS-TH-04006` x `160`
    - `RM-PET-RO-00001` x `1.952` -> `Buy`
    - `PK-BLS-PS-01014` x `160` -> `Buy`

    ## ตัวอย่างชั้นลึก: FG-PSS-TH-04001
    - `SM-PSS-TH-02001` x `160` -> Plastic branch
    - `SM-PSS-TH-02002` x `160` -> Plastic branch
    - `SM-JOI-PK-02002` x `160` -> Plastic branch
    - `RM-FIL-PS-01004` x `160` -> Mixed branch (Buy + Pharma + Plastic child)
    - `SO-PSS-LO-01001` x `0.176` -> Pharma branch

    ## Plastic branch
    - `SM-PSS-TH-02001` -> route `61 + 63`
    - `SM-PSS-TH-02002` -> route `61 + 63`
    - `SM-JOI-PK-02002` -> route `61 + 63`
    - Rule `145` -> `Transfer Plastic`
    - Rule `147` -> `Manufacturing Plastic`

    ## Pharma branch
    - `SO-PSS-LO-01001` -> route `62`
    - `RM-FIL-PS-01004` -> route `60 + 62 + Buy`
    - Rule `144` -> `Transfer Pharma`
    - Rule `146` -> `Manufacturing Pharma`

    ## Buy branch
    - `PK-CAR-PS-01003`
    - `PK-BOX-PS-01002`
    - `PK-SHF-PS-01003`
    - `RM-PET-RO-00001`
    - `PK-BLS-PS-01014`
    - Rule `7` -> `Buy` -> `Receipts`

    ## Mermaid
    ```mermaid
    flowchart TD
        A[Orderpoint 249\\nGMP/Stock\\nroute = Manufacture] --> B[FG-PSS-TH-01005\\nroute = Manufacture (Pharma)]
        B --> C[Rule 146\\nManufacturing Pharma]
        C --> D[FG-PSS-TH-02001 x16]
        C --> E[PK-CAR-PS-01003 x1\\nBuy]
        C --> F[เทปกาว x1\\nBuy]
        D --> G[FG-PSS-TH-03001 x160]
        D --> H[PK-BOX-PS-01002\\nBuy]
        D --> I[PK-SHF-PS-01003\\nBuy]
        G --> J[FG-PSS-TH-04001..04006]
        G --> K[RM-PET-RO-00001\\nBuy]
        G --> L[PK-BLS-PS-01014\\nBuy]
        J --> M[Plastic branch\\nSM-PSS / SM-JOI]
        J --> N[Pharma branch\\nSO-PSS / RM-FIL]
        M --> O[Rule 145 + 147\\nINT-PL + MO-PL]
        N --> P[Rule 144 + 146\\nINT-PH + MO-PH]
        E --> Q[Rule 7\\nReceipts]
        F --> Q
        H --> Q
        I --> Q
        K --> Q
        L --> Q
    ```
    """
).strip() + "\n"


def set_run_font(run, size, bold=False, color=None):
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_title_box(slide, title, subtitle=None):
    shape = slide.shapes.add_textbox(Inches(0.6), Inches(0.4), Inches(12.0), Inches(1.0))
    tf = shape.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, 24, bold=True, color=RGBColor(27, 47, 77))
    p.alignment = PP_ALIGN.LEFT
    if subtitle:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = subtitle
        set_run_font(r, 11, color=RGBColor(92, 102, 115))


def add_bullets(slide, bullets):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.6), Inches(11.8), Inches(5.2)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 248, 252)
    shape.line.color.rgb = RGBColor(208, 219, 231)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.level = 0
        p.space_after = Pt(10)
        r = p.add_run()
        r.text = bullet
        set_run_font(r, 18, color=RGBColor(38, 50, 56))


def build_pptx():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide_spec in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title_box(slide, slide_spec["title"], slide_spec.get("subtitle"))
        add_bullets(slide, slide_spec["bullets"])
    prs.save(PPTX_OUT)


def make_svg():
    boxes = [
        (40, 40, 220, 70, "#E8F1FF", "Orderpoint 249\nGMP/Stock\nroute = Manufacture"),
        (320, 40, 240, 70, "#EAF7EA", "FG-PSS-TH-01005\nroute = Manufacture (Pharma)"),
        (620, 40, 220, 70, "#FFF5E6", "Rule 146\nManufacturing Pharma"),
        (40, 170, 250, 90, "#F7F0FF", "Top BOM\nFG-PSS-TH-02001 x16\nPK-CAR-PS-01003 x1\nเทปกาว x1"),
        (340, 170, 250, 90, "#F7F0FF", "Mid BOM\nFG-PSS-TH-03001 x160\nPK-BOX-PS-01002\nPK-SHF-PS-01003"),
        (640, 170, 250, 90, "#F7F0FF", "Color layer\nFG-PSS-TH-04001..04006\nRM-PET-RO-00001\nPK-BLS-PS-01014"),
        (40, 330, 260, 120, "#E8F1FF", "Plastic branch\nSM-PSS-TH-02001\nSM-PSS-TH-02002\nSM-JOI-PK-02002\nroute 61 + 63"),
        (340, 330, 260, 120, "#EAF7EA", "Pharma branch\nSO-PSS-LO-01001\nRM-FIL-PS-01004\nroute 60 + 62"),
        (640, 330, 250, 120, "#FFF5E6", "Buy branch\nPK-CAR / PK-BOX / PK-SHF\nRM-PET / PK-BLS\nroute 5"),
        (60, 500, 220, 80, "#E8F1FF", "Transfer Plastic\nINT-PL\nRule 145"),
        (340, 500, 240, 80, "#EAF7EA", "Transfer / MO Pharma\nINT-PH / MO-PH\nRule 144 + 146"),
        (660, 500, 200, 80, "#FFF5E6", "Receipts\nRule 7"),
    ]
    lines = [
        (260, 75, 320, 75),
        (560, 75, 620, 75),
        (730, 110, 165, 170),
        (730, 110, 465, 170),
        (730, 110, 765, 170),
        (165, 260, 465, 260),
        (465, 260, 765, 260),
        (765, 260, 170, 330),
        (765, 260, 470, 330),
        (765, 260, 765, 330),
        (170, 450, 170, 500),
        (470, 450, 460, 500),
        (765, 450, 765, 500),
    ]
    parts = [
        '<svg xmlns="http://www.w3.org/2000/svg" width="960" height="640" viewBox="0 0 960 640">',
        '<rect width="100%" height="100%" fill="#ffffff"/>',
        '<text x="40" y="22" font-family="Tahoma" font-size="20" font-weight="700" fill="#1B2F4D">Flow Diagram: FG-PSS-TH-01005 (UAT)</text>',
    ]
    for x1, y1, x2, y2 in lines:
        parts.append(
            f'<line x1="{x1}" y1="{y1}" x2="{x2}" y2="{y2}" stroke="#7A8CA5" stroke-width="2" marker-end="url(#arrow)"/>'
        )
    parts.insert(
        1,
        '<defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="8" refY="3" orient="auto" markerUnits="strokeWidth"><path d="M0,0 L0,6 L9,3 z" fill="#7A8CA5"/></marker></defs>',
    )
    for x, y, w, h, fill, text in boxes:
        parts.append(
            f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="12" ry="12" fill="{fill}" stroke="#B7C6D8" stroke-width="1.5"/>'
        )
        for idx, line in enumerate(text.split("\n")):
            parts.append(
                f'<text x="{x + 12}" y="{y + 24 + idx * 20}" font-family="Tahoma" font-size="15" fill="#263238">{escape(line)}</text>'
            )
    parts.append("</svg>")
    FLOW_SVG_OUT.write_text("\n".join(parts), encoding="utf-8")


def main():
    SLIDE_MD_OUT.write_text(SLIDE_MD, encoding="utf-8")
    FLOW_MD_OUT.write_text(FLOW_MD, encoding="utf-8")
    TRAINER_OUT.write_text(TRAINER_MD, encoding="utf-8")
    make_svg()
    build_pptx()
    print(SLIDE_MD_OUT)
    print(FLOW_MD_OUT)
    print(FLOW_SVG_OUT)
    print(TRAINER_OUT)
    print(PPTX_OUT)


if __name__ == "__main__":
    main()
