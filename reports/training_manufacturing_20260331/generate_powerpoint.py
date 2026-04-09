from pathlib import Path

from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor


OUT = Path(r"c:\365_project\TheCool18e\Dev\reports\training_manufacturing_20260331\manufacturing_training_20260331.pptx")


SLIDES = [
    {
        "layout": "title",
        "title": "Manufacturing Training",
        "subtitle": "Odoo 18 Enterprise\nPlastic + Pharma Integrated Flow\nSession 3 ชั่วโมง",
    },
    {
        "title": "Session Objective",
        "bullets": [
            "เข้าใจภาพรวมระบบ Manufacturing ของบริษัท",
            "เข้าใจ Promotion, MTO, MTS / Min-Max",
            "เข้าใจ Shopfloor, Scrap และ Mold",
            "รู้ว่าระบบ auto อะไร และต้องเช็กอะไร",
        ],
    },
    {
        "title": "Business Picture",
        "bullets": [
            "โรงงานมี 2 ฝั่ง: Plastic และ Pharma",
            "Semi หลายตัวผลิตที่ Plastic แล้วส่งต่อไป Pharma",
            "Demand 1 ตัว อาจแตกเป็นหลาย MO, Transfer และ PO",
            "เป้าหมายคือให้ระบบ auto chain ให้มากที่สุด",
        ],
    },
    {
        "title": "Document Map",
        "bullets": [
            "SO / Quotation",
            "Replenishment / Orderpoint",
            "Manufacturing Order / Work Order",
            "Transfer / Delivery / Invoice / Journal Entry",
        ],
    },
    {
        "title": "Settings That Matter",
        "bullets": [
            "Product Routes",
            "Manufacturing Type",
            "BoM และ BoM Picking Type",
            "Orderpoint และ Vendor",
            "Workcenter, Operation Type, Mold Mapping",
        ],
    },
    {
        "title": "Factory Separation Logic",
        "bullets": [
            "Manufacturing Plastic",
            "Manufacturing Pharma",
            "Transfer Plastic",
            "Transfer Pharma",
            "แยก operation type ให้ถูกก่อน flow จะ auto ถูก",
        ],
    },
    {
        "title": "Flow 1: Promotion / SO 0 บาท",
        "bullets": [
            "SO free item และ FOC line",
            "SO -> MO -> Transfer -> Delivery -> Invoice",
            "FOC ไม่ได้แปลว่าไม่มี stock movement",
            "FOC ไม่ได้แปลว่าไม่มี accounting",
        ],
        "footer": "Demo: S11563, SOB-263069, GMP/OUT/02504, INV-D/26/03/00960",
    },
    {
        "title": "Flow 2: SO -> MTO",
        "bullets": [
            "Demand มาจาก sale โดยตรง",
            "ระบบแตก child MO ตาม BoM chain",
            "ตรวจ procurement group และ smart buttons",
            "แยก make_to_order กับ delivery rule ให้ชัด",
        ],
        "footer": "Demo: SOB-263070, GMP/PICK/03572, GMP/OUT/02505, INV-D/26/03/00961",
    },
    {
        "title": "Flow 3: MTS / Min-Max",
        "bullets": [
            "ไม่มี SO ก็ผลิตได้จาก Reordering Rule",
            "ของพอ ระบบไม่สร้าง procurement เพิ่ม",
            "ของไม่พอ ระบบตัดสิน route ไป MO หรือ PO",
            "0/0 orderpoint ทำงานแบบเติม shortage กลับมาที่ 0",
        ],
        "footer": "Demo: FG-PNC-TH-01001, GMP/MOPH/00011",
    },
    {
        "title": "Example FG-PSS-TH-01005",
        "bullets": [
            "ปัจจุบันวิ่งจาก Min/Max ไม่ใช่ MTO จาก SO",
            "MO แม่อยู่ฝั่ง Manufacturing Pharma",
            "แตก child chain ลงทั้ง FG, Semi, Solution",
            "ใต้ chain มีทั้ง Plastic และ Pharma",
        ],
    },
    {
        "title": "Multi-level Manufacturing",
        "bullets": [
            "FG carton layer",
            "FG 5 โหล layer",
            "FG สี layer",
            "Semi plastic layer",
            "Solution layer",
        ],
    },
    {
        "title": "Shopfloor Execution",
        "bullets": [
            "Start, Good Qty, Reject Qty, Done",
            "Workorder คือหน้างานจริงของ operator",
            "Scrap มีผลต่อ stock และ cost",
            "กรอก qty ผิด จะกระทบ valuation ทันที",
        ],
    },
    {
        "title": "Parallel Workcenter Rules",
        "bullets": [
            "Done หรือ Cancel ไม่เข้ากระจาย planned qty แล้ว",
            "Cancelled Workorder ไม่แสดงใน Shopfloor",
            "Logic นี้ช่วยกัน qty เพี้ยนในงาน parallel",
            "Supervisor ต้องดูเฉพาะ active workorder",
        ],
    },
    {
        "title": "Mold Management",
        "bullets": [
            "Mold map กับ product และ workcenter",
            "ระบบ auto assign mold ตาม matrix",
            "Mold เดียวไม่ถูก assign ซ้ำข้าม parallel WO",
            "Mold life ขึ้นจาก output จริง",
        ],
        "footer": "Demo: GMP/MOPL/00014, SP/00011",
    },
    {
        "title": "Transfer Logic",
        "bullets": [
            "Semi จาก Plastic ไป Pharma ผ่าน internal transfer",
            "แยก Transfer Plastic และ Transfer Pharma",
            "เช็กจาก Operation Type ไม่ใช่ชื่อสินค้าอย่างเดียว",
            "ถ้า transfer ไม่มา ให้เช็ก route และ location",
        ],
    },
    {
        "title": "Purchasing Logic",
        "bullets": [
            "ขาดของแล้วจะไป PO ก็ต่อเมื่อ route เป็น Buy",
            "และต้องมี Vendor พร้อม",
            "ถ้า orderpoint route เป็น Manufacture ของจะไป MO ก่อน",
            "Buy route มีแต่ไม่มี vendor ระบบไปต่อไม่สุด",
        ],
    },
    {
        "title": "Costing Logic",
        "bullets": [
            "Actual cost แยกเป็น Raw, Machine, Labor, Mold",
            "Std cost กับ Actual cost คนละมุมมอง",
            "Production และ Accounting ต้องเห็นภาพเดียวกัน",
            "อย่าดูเฉพาะ invoice ให้ดู stock valuation ด้วย",
        ],
    },
    {
        "title": "Common Failure Points",
        "bullets": [
            "Route ผิดฝั่ง",
            "Vendor ไม่มี",
            "Orderpoint ไม่ครบ",
            "Mold ไม่ map",
            "Stock อยู่ผิด location",
            "Transfer ไม่เกิด",
        ],
    },
    {
        "title": "Demo Documents",
        "bullets": [
            "Promotion quote: S11563",
            "FOC flow: SOB-263069",
            "MTO flow: SOB-263070",
            "MTS flow: GMP/MOPH/00011",
            "100000 production trace: GMP/MOPH/00001",
            "Mold and shopfloor: GMP/MOPL/00014",
        ],
    },
    {
        "title": "UAT Checklist",
        "bullets": [
            "Promotion",
            "MTO",
            "MTS",
            "Shopfloor",
            "Mold",
            "Scrap",
            "Costing",
        ],
    },
    {
        "title": "Closing",
        "bullets": [
            "ถ้าของหมด ระบบจะไป MO หรือ PO เพราะอะไร",
            "อะไรเป็น auto และอะไรต้องตัดสินใจเอง",
            "ต้องเช็ก route, BoM, orderpoint, vendor, mold mapping ก่อนเสมอ",
        ],
    },
]


def set_run_font(run, size, bold=False, color=None):
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(6.75), Inches(12.3), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, 10, color=RGBColor(102, 102, 102))


def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)

    title_shape = slide.shapes.title
    title_shape.text = title
    for p in title_shape.text_frame.paragraphs:
        for run in p.runs:
            set_run_font(run, 28, bold=True, color=RGBColor(31, 41, 55))

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    for p in subtitle_shape.text_frame.paragraphs:
        for run in p.runs:
            set_run_font(run, 16, color=RGBColor(75, 85, 99))

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.45))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(22, 101, 52)
    band.line.fill.background()


def add_content_slide(prs, title, bullets, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_shape = slide.shapes.title
    title_shape.text = title
    for p in title_shape.text_frame.paragraphs:
        for run in p.runs:
            set_run_font(run, 24, bold=True, color=RGBColor(22, 101, 52))

    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            set_run_font(run, 18, color=RGBColor(31, 41, 55))

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.8), Inches(0.18), Inches(5.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(22, 101, 52)
    accent.line.fill.background()

    if footer:
        add_footer(slide, footer)


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        if slide_data.get("layout") == "title":
            add_title_slide(prs, slide_data["title"], slide_data["subtitle"])
        else:
            add_content_slide(
                prs,
                slide_data["title"],
                slide_data["bullets"],
                slide_data.get("footer"),
            )

    prs.save(OUT)
    print(f"saved {OUT}")


if __name__ == "__main__":
    main()
