from pathlib import Path
from textwrap import fill

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.util import Inches, Pt


BASE = Path(r"C:\365_project\TheCool18e\Dev\reports\training_manufacturing_20260331")
SHOT_DIR = BASE / "uat_screenshots"
OUT_PPTX = BASE / "manufacturing_training_uat_walkthrough_20260331.pptx"
OUT_TXT = BASE / "manufacturing_training_uat_walkthrough_speaker_20260331.txt"


GREEN = RGBColor(22, 101, 52)
GREEN_DARK = RGBColor(15, 70, 38)
TEXT = RGBColor(31, 41, 55)
MUTED = RGBColor(75, 85, 99)
BG = RGBColor(245, 247, 250)
WHITE = RGBColor(255, 255, 255)
ACCENT = RGBColor(222, 243, 230)
LINE = RGBColor(210, 214, 220)


SLIDES = [
    {
        "kind": "title",
        "title": "Manufacturing Training Walk-through",
        "subtitle": "UAT Database\nOdoo 18 Enterprise\nPlastic + Pharma + PO + Mold + Costing",
        "speaker": [
            "เปิดคลาสด้วยการบอกว่าชุดนี้ใช้ภาพจริงจากฐาน UAT ไม่ใช่ mock screen",
            "บอกผู้เข้าอบรมว่าวันนี้จะไล่ตั้งแต่ master data ไปจนถึง transaction ปลายทาง",
        ],
    },
    {
        "kind": "agenda",
        "title": "Training Agenda",
        "items": [
            "1. ภาพรวมระบบและแผนที่เอกสาร",
            "2. Settings และ Master Data ที่ทำให้ระบบ auto",
            "3. Promotion / SO 0 บาท และ free item",
            "4. SO -> MTO",
            "5. MTS / Min-Max และ Replenishment",
            "6. PO จาก shortage และ Buy route",
            "7. Work Order / Shop Floor / Scrap",
            "8. Workcenter / Mold Mapping",
            "9. Transfer Plastic / Pharma",
            "10. Invoice / Costing / Closing Checklist",
        ],
        "speaker": [
            "สไลด์นี้ใช้เปิด session ว่าวันนี้จะสอนจาก setting ไปถึง transaction ปลายทาง",
            "ย้ำให้ผู้เข้าอบรมเห็นว่าไม่ได้เรียนแค่ Manufacturing app แต่จะเชื่อม Sales, Inventory, Purchase และ Accounting ด้วย",
            "หัวข้อ PO ถูกแทรกไว้ชัดเจน เพราะ shortage ไม่ได้จบที่ MO ทุกกรณี",
        ],
    },
]

SLIDES.extend(
    [
        {
            "title": "Home Dashboard",
            "section": "Opening",
            "menu": "Home",
            "screenshot": "01_home_dashboard.png",
            "what": [
                "ใช้เป็นจุดเริ่มต้นเพื่ออธิบายว่า flow จริงไม่ได้อยู่ในแอปเดียว",
                "ให้ผู้ใช้งานเห็นกลุ่มแอปหลักที่เกี่ยวข้องกับงานผลิต",
            ],
            "fields": [
                "`Sales` จุดเริ่มต้นของ demand จากลูกค้า",
                "`Inventory` ใช้ดู stock, transfer, reordering rule และ scrap",
                "`Manufacturing` ใช้ดู MO, WO, BoM และ workcenter",
                "`Purchase` ใช้ดู RFQ/PO เมื่อ shortage ต้องซื้อ",
                "`Accounting` ใช้ดู invoice และผลทางบัญชี",
                "`GMP Shop Floor` เป็นหน้าหน้างานสำหรับการปฏิบัติการจริง",
            ],
            "speaker": [
                "เริ่มจากหน้า Home เพื่อให้ทุกคนเห็น landscape ของระบบก่อน",
                "ย้ำว่าการผลิตของบริษัทไม่ได้จบใน Manufacturing app แต่เชื่อมกับ Sales, Inventory, Purchase และ Accounting",
                "บอกผู้เข้าอบรมว่าระหว่าง training จะวนกลับมาที่หน้าหลักนี้เสมอเมื่ออธิบาย flow ข้ามแอป",
            ],
        },
        {
            "title": "Manufacturing Overview",
            "section": "Opening",
            "menu": "Manufacturing > Overview",
            "screenshot": "02_manufacturing_overview.png",
            "what": [
                "ใช้สอนโครงเมนูของ Manufacturing และหน้าภาพรวม work centers",
                "เป็นหน้าที่ supervisor ใช้ดูภาระงานและดูว่ามี work orders อยู่ที่เครื่องไหน",
            ],
            "fields": [
                "`Overview` ใช้ดูสถานะ workcenter เป็นภาพรวม",
                "`Operations` รวม MO, WO และ planning ที่ต้องติดตาม",
                "`Products` รวม product master และ BoM",
                "`Configuration` ใช้ตั้ง workcenter, operation type และ master data ที่มีผลกับ flow",
                "`PLAN ORDERS` หรือ `WORK ORDERS` บอกว่ามีงานที่เกี่ยวข้องกับ workcenter นั้นหรือไม่",
            ],
            "speaker": [
                "ชี้ให้ทีมเห็นว่า Manufacturing app มีทั้งส่วนปฏิบัติการและส่วนตั้งค่า",
                "อธิบายว่า Overview เหมาะกับหัวหน้างานและ planner เพราะเห็นเครื่องจำนวนมากในหน้าเดียว",
                "ใช้สไลด์นี้เชื่อมไปสู่เรื่อง workcenter, work order และ mold later in the session",
            ],
        },
        {
            "title": "Product Master: FG-PSS-TH-01005",
            "section": "Master Data",
            "menu": "Manufacturing > Products > Products",
            "screenshot": "03_product_fg_pss_th_01005.png",
            "what": [
                "ใช้สอนว่าหน้าสินค้าเป็นต้นทางของ logic หลายอย่างในระบบ",
                "ตัวอย่างนี้เป็น finished good ฝั่ง Pharma ที่ถูกใช้ตลอด training",
            ],
            "fields": [
                "`Product Type` ระบุว่าสินค้าเป็น Goods, Service หรือ Combo",
                "`Manufacturing Type` ใช้แยก Plastic / Pharma / Packaging เพื่อพา flow ไป operation type ที่ถูก",
                "`Track Inventory` กำหนดรูปแบบการติดตาม lot/serial/quantity",
                "`Invoicing Policy` บอกว่าสร้าง invoice ตาม ordered หรือ delivered quantities",
                "`Sales Price` คือราคาขายมาตรฐาน",
                "`Cost` คือต้นทุนอ้างอิงที่ใช้ใน master data",
                "`Category` และ `Reference` ช่วยให้ route, บัญชี และ reporting ทำงานถูก",
            ],
            "speaker": [
                "อธิบายว่าถ้าตั้ง Product master ผิด ระบบจะ auto ผิดทั้ง chain",
                "ย้ำว่า Manufacturing Type ไม่ใช่แค่ label แต่เป็นตัวคุมว่าเอกสารจะไป Plastic หรือ Pharma",
                "ชี้ smart buttons ด้านบน เช่น On Hand, Forecasted และ Bill of Materials เพื่อให้ผู้ใช้รู้ว่าจะ drill down จากตรงไหน",
            ],
        },
        {
            "title": "Product Inventory & Routes",
            "section": "Master Data",
            "menu": "Manufacturing > Products > Products > Inventory",
            "screenshot": "03b_product_fg_pss_th_01005_inventory_routes.png",
            "what": [
                "ใช้สอน field ในแท็บ Inventory ที่ทำให้ระบบตัดสินใจว่าจะไป MO หรือ PO",
                "สไลด์นี้คือหัวใจของการอธิบาย route และ replenishment",
            ],
            "fields": [
                "`Routes` คือ route จริงที่ product ใช้ เช่น Manufacture (Pharma), Buy, Replenish on Order",
                "`Category Routes` คือ route ที่ inherited มาจาก category หรือ warehouse",
                "`View Diagram` ใช้ไล่ logic route แบบภาพ",
                "`Responsible` บอกคนรับผิดชอบด้าน logistics/master data",
                "`Customer Lead Time` และ `MFG Lead Time` ใช้ประกอบการวางแผนเวลา",
                "`Expiration / Best Before / Removal / Alert Date` สำคัญกับสินค้าที่มีการควบคุม lot และ shelf life",
            ],
            "speaker": [
                "สไลด์นี้ต้องสอนช้าและชัด เพราะเป็นจุดที่ตอบคำถามยอดฮิตว่า ของขาดแล้วทำไมไป MO หรือ PO",
                "ชี้ route ที่ถูกติ๊กใน UAT ว่าตัวนี้ใช้ Manufacture (Pharma) ไม่ได้ใช้ MTO จาก sale โดยตรง",
                "อธิบายความต่างระหว่าง route ที่เลือกใน product กับ category routes ที่ inherited มา",
            ],
        },
        {
            "title": "Reordering Rule / Min-Max",
            "section": "Planning",
            "menu": "Inventory > Configuration > Reordering Rules",
            "screenshot": "04_reordering_rule_fg_pss_th_01005.png",
            "what": [
                "ใช้สอน flow MTS ว่าระบบเติม stock อย่างไรเมื่อไม่ได้เริ่มจาก SO",
                "ตัวอย่างนี้เป็น orderpoint ของ FG-PSS-TH-01005 ที่ใช้กระตุ้น replenishment",
            ],
            "fields": [
                "`Product` คือสินค้าที่ reordering rule ควบคุม",
                "`Warehouse` และ `Location` ระบุว่าระบบจะดู shortage ที่ stock ไหน",
                "`Min Quantity` คือระดับต่ำสุดที่ยอมให้มีใน location",
                "`Max Quantity` คือระดับเป้าหมายที่ระบบจะเติมกลับไป",
                "`Quantity Multiple` ใช้บังคับปริมาณการเติมเป็น lot",
                "`Forecast Description` ใช้ drill down เพื่อดูว่าความต้องการมาจากเอกสารใด",
            ],
            "speaker": [
                "อธิบายว่า min/max ไม่ได้หมายความว่าต้องมีตัวเลขสูงเสมอไป บางกรณีใช้ 0/0 เพื่อเติม shortage กลับมาที่ศูนย์",
                "เชื่อมกับสไลด์ product routes ว่า route ของ orderpoint จะบอกต่อว่าขาดแล้วไป Manufacture หรือ Buy",
                "ย้ำว่าถ้าไม่มี reordering rule ระบบ MTS จะไม่ auto สร้าง procurement",
            ],
        },
        {
            "title": "Bill of Materials",
            "section": "Master Data",
            "menu": "Manufacturing > Products > Bills of Materials",
            "screenshot": "05_bom_fg_pss_th_01005.png",
            "what": [
                "ใช้สอนว่า MO ตัวแม่เรียก component อะไรบ้างในระดับแรก",
                "สไลด์นี้ช่วยปูพื้นก่อนอธิบาย child MO และ transfer",
            ],
            "fields": [
                "`Product` คือสินค้าที่ BoM นี้อธิบาย",
                "`Quantity` คือ standard batch ของ BoM",
                "`BoM Type` ระบุว่าเป็น Manufacture, Kit หรือ Subcontracting",
                "`Components` คือรายการวัตถุดิบ/กึ่งสำเร็จรูปที่ต้องใช้",
                "`Operations` ใช้ผูกขั้นตอนการผลิตกับ workcenter หรือ operation template",
                "`Miscellaneous` มักใช้ดู operation type หรือ option เฉพาะของ BoM",
            ],
            "speaker": [
                "อธิบายให้ผู้เรียนเห็นว่า BoM นี้ยังเป็นแค่ชั้นแรกของ FG ไม่ได้หมายความว่าทุกอย่างซื้อหรือผลิตในชั้นเดียว",
                "ชี้ component สำคัญ เช่น FG ย่อย, packaging และเทปกาว เพื่อปูไปเรื่อง MO กับ PO",
                "ถ้าต้องการ trace chain ลึก ต้องกดเข้าดู child products ต่อ ไม่ใช่ดู BoM เดียวแล้วสรุปทันที",
            ],
        },
        {
            "title": "Purchase Order from Shortage",
            "section": "Planning",
            "menu": "Purchase > Orders > Requests for Quotation",
            "screenshot": "09_po_p00014.png",
            "what": [
                "ใช้สอนว่าขาดของแล้วระบบไม่ได้ไป MO ทุกครั้ง",
                "ตัวอย่างนี้เป็น RFQ จริงสำหรับวัตถุดิบและ packaging ที่ route เป็น Buy",
            ],
            "fields": [
                "`Vendor` คือคู่ค้าที่ระบบจะใช้สร้าง RFQ/PO",
                "`Order Deadline` และ `Expected Arrival` ใช้ประกอบการวางแผนรับของ",
                "`Deliver To` ระบุปลายทางการรับเข้า",
                "`Products` คือรายการที่ถูกสั่งซื้อจาก shortage หรือการวางแผน",
                "`Quantity`, `UoM`, `Unit Price`, `Taxes` ใช้คำนวณมูลค่าการซื้อ",
                "`Price Over Limit` เป็นสัญญาณเตือนให้ตรวจราคาก่อน confirm",
            ],
            "speaker": [
                "ย้ำว่า PO จะเกิดได้ต่อเมื่อสินค้านั้นมี Buy route และมี vendor พร้อม",
                "ชี้ว่าหน้า PO เป็นจุดที่ Purchase เข้ามารับช่วงจาก planning หรือ shortage logic",
                "ใช้สไลด์นี้ตอบคำถามว่า ถ้าไม่มี vendor แม้ route ถูก ระบบก็ไปต่อ PO ไม่สุด",
            ],
        },
    ]
)

SLIDES.extend(
    [
        {
            "title": "Mold Master",
            "section": "Mold",
            "menu": "Manufacturing > Configuration > Work Centers (Mold)",
            "screenshot": "11_mold_upper_w01_general.png",
            "what": [
                "ใช้สอน master data ของ mold ที่ตอนนี้เปิดให้เห็น field Is Mold? แล้ว",
                "เป็นหน้าอ้างอิงเรื่อง mold cost, mold life และ cavities",
            ],
            "fields": [
                "`Is Mold?` ต้องติ๊กเพื่อให้ record นี้ทำงานใน logic mold",
                "`Mold Cost / Hour` ใช้คำนวณต้นทุน mold",
                "`Cavities` ใช้ตีความ output ต่อหนึ่ง shot",
                "`Mold Life Limit (Shots)` คืออายุใช้งานสูงสุดตามแผน",
                "`Mold Status` บอกว่าตอนนี้อยู่ในระดับ Normal, Warning หรือ Full",
                "`Shots Taken` smart button ด้านบนใช้ดูการใช้งานสะสมเทียบ limit",
            ],
            "speaker": [
                "อธิบายว่าหน้า mold ใช้ model เดียวกับ workcenter แต่ความหมายของ field ต่างกัน",
                "ย้ำว่าถ้าไม่ได้ติ๊ก Is Mold ระบบจะไม่ treat record นี้เป็น mold และ matrix/logic บางอย่างจะไม่ทำงาน",
                "ใช้สไลด์นี้อธิบายความสัมพันธ์ระหว่าง mold cost, mold life และการบำรุงรักษา",
            ],
        },
        {
            "title": "Mold Matrix",
            "section": "Mold",
            "menu": "Manufacturing > Configuration > Work Centers (Mold) > Mold Matrix",
            "screenshot": "11b_mold_upper_w01_matrix.png",
            "what": [
                "ใช้สอนสองเรื่องพร้อมกันคือ compatible machines และ produced products efficiency",
                "เป็นข้อมูลฝั่ง mold ที่ช่วยให้ระบบรู้ว่า mold นี้ใช้กับสินค้าอะไรและเครื่องไหนได้",
            ],
            "fields": [
                "`Compatible Machines` คือรายการเครื่องที่รับ mold ตัวนี้ได้",
                "`Machines are linked back...` อธิบายว่าความสัมพันธ์นี้มาจาก machine side ได้ด้วย",
                "`Product` คือสินค้าที่ mold นี้ผลิตได้",
                "`Cycle Time (s)` คือเวลามาตรฐานต่อรอบ",
                "`Units / Hour` ใช้ช่วยคำนวณความสามารถในการผลิต",
                "`Add a line` ใช้เพิ่ม product efficiency row ใหม่",
            ],
            "speaker": [
                "ย้ำว่าถ้าขาดฝั่ง machine หรือฝั่ง product ใดฝั่งหนึ่ง ระบบ auto assign mold จะไม่สมบูรณ์",
                "สอนให้ทีมอ่านว่าตารางนี้คือสะพานระหว่าง master data กับ execution จริง",
                "เชื่อมกับ issue ที่เราเคยแก้ว่า mold เดียวไม่ควรถูก assign ซ้ำข้าม parallel workorders",
            ],
        },
        {
            "title": "Transfer Plastic",
            "section": "Inventory",
            "menu": "Inventory > Operations > Transfers",
            "screenshot": "12_transfer_plastic_gmp_trpl_00006.png",
            "what": [
                "ใช้สอน internal transfer ฝั่ง Plastic และการส่ง semi ต่อไปยัง stock/ขั้นถัดไป",
                "เป็นหลักฐานว่าระบบแยก operation type ของ Plastic ชัดเจน",
            ],
            "fields": [
                "`Operation Type` ต้องเป็น Transfer Plastic",
                "`Source Location` และ `Destination Location` ระบุทิศทางการย้าย stock",
                "`Source Document` ใช้ trace กลับไป MO ที่เป็นต้นทาง",
                "`Product` ใน operation lines คือ semi หรือชิ้นส่วนที่ย้ายจริง",
                "`Demand` กับ `Quantity` ใช้ดูความต้องการเทียบกับ qty done/processed",
                "`Validate` คือขั้นตอนยืนยัน movement เข้าสต็อกปลายทาง",
            ],
            "speaker": [
                "สอนให้ดู operation type ก่อนเสมอ ไม่ใช้ชื่อสินค้าอย่างเดียวในการบอกว่าฝั่งไหน",
                "อธิบายว่าเอกสาร transfer คือสะพานเชื่อมระหว่างการผลิตกับ stock movement",
                "หาก transfer ไม่เกิด ให้ย้อนกลับไปเช็ก routes, BoM picking type และ source/destination logic",
            ],
        },
        {
            "title": "Transfer Pharma",
            "section": "Inventory",
            "menu": "Inventory > Operations > Transfers",
            "screenshot": "13_transfer_pharma_gmp_trph_00001.png",
            "what": [
                "ใช้สอน internal transfer ฝั่ง Pharma และการรับชิ้นส่วน/FG ที่เกี่ยวข้องกับโรงงานยา",
                "จับคู่กับสไลด์ก่อนหน้าเพื่อเทียบ Plastic กับ Pharma แบบเห็นภาพ",
            ],
            "fields": [
                "`Operation Type` ต้องเป็น Transfer Pharma",
                "`Source Document` อาจรวม MO หลายใบที่ feed เข้ามาใน transfer เดียว",
                "`Operations` table ใช้ดูสินค้าแต่ละตัวที่ถูกรวมมาจากงานผลิต",
                "`Demand` ช่วยอ่าน requirement ของการเคลื่อนย้าย",
                "`Ready / Done` status ช่วยบอกว่าการเคลื่อนย้ายถึงขั้นไหนแล้ว",
            ],
            "speaker": [
                "ใช้สไลด์นี้เน้นว่าระบบสามารถแยก Plastic กับ Pharma ในระดับ transfer ได้จริง",
                "ให้ผู้เรียนสังเกตว่า source document อาจมีหลาย MO รวมอยู่ได้ในเอกสารเดียว",
                "อธิบายว่า inventory team ต้องอ่าน transfer ให้เป็นเพื่อ trace ปัญหา stock ข้ามโรงงาน",
            ],
        },
        {
            "title": "Invoice & Accounting Endpoint",
            "section": "Accounting",
            "menu": "Accounting > Customers > Invoices",
            "screenshot": "14_invoice_inv_d_26_04_00001.png",
            "what": [
                "ใช้สอนปลายทางด้านเอกสารบัญชีหลังจาก sale และ delivery วิ่งต่อมา",
                "แม้ training หลักอยู่ฝั่ง Manufacturing แต่ต้องจบที่ภาพรวม accounting ด้วย",
            ],
            "fields": [
                "`Customer`, `Invoice Date`, `Payment terms`, `Journal` คือ header หลักของ invoice",
                "`Invoice Lines` ใช้ดูสินค้า, quantity, price และภาษี",
                "`Account` คือบัญชีรายได้หรือบัญชีที่ผูกจาก product/accounting config",
                "`Analytic` ใช้สำหรับ reporting หรือการกระจายต้นทุน/รายได้",
                "`Untaxed Amount`, `VAT`, `Total` ใช้ยืนยันยอด",
                "`Posted` คือสถานะที่พร้อมลงบัญชีจริง",
            ],
            "speaker": [
                "ชี้ว่าฝั่ง Accounting เป็นปลายทางของ flow และเป็นจุดที่ผู้บริหารมักใช้ตรวจผลลัพธ์",
                "อธิบายให้ทีมรู้ว่า invoice ไม่ได้ยืนเดี่ยว แต่ย้อน trace กลับไป SO, delivery และ product accounts ได้",
                "ถ้าต้องการอธิบาย cost ลึกกว่านี้ ให้เสริมว่าต้นทุนจริงยังต้องดู stock valuation และ journal ที่เกิดจาก MO ด้วย",
            ],
        },
        {
            "kind": "closing",
            "title": "Closing Checklist",
            "items": [
                "เช็ก Product Routes และ Manufacturing Type ก่อนทุกครั้ง",
                "เช็ก BoM และ child chain เมื่อ MO ไม่แตกต่อ",
                "เช็ก Reordering Rule, vendor และ Buy route เมื่อของขาดแล้ว PO ไม่มา",
                "เช็ก Workcenter / Mold Matrix เมื่อหน้างานเลือกเครื่องหรือ mold ไม่ได้",
                "เช็ก Transfer Operation Type เมื่อ stock ข้ามโรงผิดฝั่ง",
                "เช็ก Invoice, Valuation และ Journal เมื่อต้อง trace ผลทางบัญชี",
            ],
            "speaker": [
                "ปิด training ด้วย checklist ที่ผู้ใช้งานต้องจำจริง",
                "ย้ำว่าก่อนแจ้งปัญหา ควรเช็ก 5 จุดหลัก: route, BoM, orderpoint, vendor, mold mapping",
                "ถ้าทีมตอบได้ว่าขาดของแล้วจะไป MO หรือ PO เพราะอะไร แปลว่าจับระบบได้แล้ว",
            ],
        },
    ]
)


def set_run_style(run, size, bold=False, color=TEXT):
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_text(paragraph, text, size, bold=False, color=TEXT):
    paragraph.clear()
    run = paragraph.add_run()
    run.text = text
    set_run_style(run, size=size, bold=bold, color=color)


def add_banner(slide, title, subtitle=None):
    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.65),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.35), Inches(0.12), Inches(8.5), Inches(0.3))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    set_text(title_tf.paragraphs[0], title, 24, True, WHITE)

    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(9.4), Inches(0.14), Inches(3.25), Inches(0.25))
        sub_tf = sub_box.text_frame
        sub_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        set_text(sub_tf.paragraphs[0], subtitle, 11, False, WHITE)


def add_card(slide, left, top, width, height, title, lines):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.24))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    set_text(title_tf.paragraphs[0], title, 13, True, GREEN_DARK)

    body_box = slide.shapes.add_textbox(left + Inches(0.14), top + Inches(0.38), width - Inches(0.28), height - Inches(0.48))
    body_tf = body_box.text_frame
    body_tf.word_wrap = True
    body_tf.margin_left = 0
    body_tf.margin_right = 0
    body_tf.margin_top = 0
    body_tf.margin_bottom = 0
    body_tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    first = True
    for line in lines:
        para = body_tf.paragraphs[0] if first else body_tf.add_paragraph()
        first = False
        para.level = 0
        para.bullet = True
        para.space_after = Pt(2)
        para.line_spacing = 1.05
        run = para.add_run()
        run.text = fill(line, 36)
        set_run_style(run, 9.5, False, TEXT)


def add_screenshot(slide, filename):
    path = SHOT_DIR / filename
    frame = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(6.2),
        Inches(1.0),
        Inches(6.7),
        Inches(5.9),
    )
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    frame.line.width = Pt(1)

    if not path.exists():
        missing = slide.shapes.add_textbox(Inches(6.45), Inches(3.4), Inches(6.2), Inches(0.5))
        tf = missing.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        set_text(tf.paragraphs[0], f"Missing screenshot: {filename}", 14, True, GREEN_DARK)
        return

    slide.shapes.add_picture(str(path), Inches(6.32), Inches(1.12), width=Inches(6.46), height=Inches(5.66))


def add_footer(slide, section, menu):
    footer = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(7.12),
        Inches(13.333),
        Inches(0.38),
    )
    footer.fill.solid()
    footer.fill.fore_color.rgb = GREEN_DARK
    footer.line.fill.background()

    left_box = slide.shapes.add_textbox(Inches(0.3), Inches(7.18), Inches(3.2), Inches(0.18))
    left_tf = left_box.text_frame
    set_text(left_tf.paragraphs[0], section, 10, True, WHITE)

    right_box = slide.shapes.add_textbox(Inches(3.6), Inches(7.18), Inches(9.3), Inches(0.18))
    right_tf = right_box.text_frame
    right_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_text(right_tf.paragraphs[0], menu, 10, False, WHITE)


def build_title(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG

    hero = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.65),
        Inches(1.0),
        Inches(12.05),
        Inches(5.4),
    )
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.color.rgb = LINE

    banner = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.65),
        Inches(1.0),
        Inches(12.05),
        Inches(0.8),
    )
    banner.fill.solid()
    banner.fill.fore_color.rgb = GREEN
    banner.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10.7), Inches(1.1))
    title_tf = title_box.text_frame
    title_tf.word_wrap = True
    set_text(title_tf.paragraphs[0], data["title"], 28, True, GREEN_DARK)

    sub_box = slide.shapes.add_textbox(Inches(1.0), Inches(3.0), Inches(9.0), Inches(1.3))
    sub_tf = sub_box.text_frame
    first = True
    for line in data["subtitle"].splitlines():
        para = sub_tf.paragraphs[0] if first else sub_tf.add_paragraph()
        first = False
        set_text(para, line, 17 if first else 15, False, MUTED)

    tag = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(5.2),
        Inches(3.3),
        Inches(0.55),
    )
    tag.fill.solid()
    tag.fill.fore_color.rgb = ACCENT
    tag.line.fill.background()
    tag_box = slide.shapes.add_textbox(Inches(1.16), Inches(5.34), Inches(3.0), Inches(0.2))
    tag_tf = tag_box.text_frame
    set_text(tag_tf.paragraphs[0], "Real screenshots from UAT database", 11, True, GREEN_DARK)


def build_agenda(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, data["title"], "3-hour training flow")

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(1.0),
        Inches(11.85),
        Inches(5.95),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE

    body = slide.shapes.add_textbox(Inches(1.15), Inches(1.35), Inches(11.0), Inches(5.1))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for idx, item in enumerate(data["items"]):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.space_after = Pt(10)
        para.line_spacing = 1.1
        run = para.add_run()
        run.text = item
        set_run_style(run, 20 if idx == 0 else 18, idx == 0, GREEN_DARK if idx == 0 else TEXT)

    add_footer(slide, "Agenda", "Training sequence")


def build_content(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, data["title"], data.get("section"))
    add_card(slide, Inches(0.4), Inches(1.05), Inches(5.45), Inches(1.75), "หน้าที่ของหน้า", data["what"])
    add_card(slide, Inches(0.4), Inches(2.95), Inches(5.45), Inches(3.8), "Field ที่ต้องอธิบาย", data["fields"])
    add_screenshot(slide, data["screenshot"])
    add_footer(slide, data["section"], data["menu"])


def build_closing(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, data["title"], "Key points before go-live")

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.75),
        Inches(1.0),
        Inches(11.85),
        Inches(5.95),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE

    body = slide.shapes.add_textbox(Inches(1.1), Inches(1.35), Inches(11.1), Inches(5.15))
    tf = body.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0

    for idx, item in enumerate(data["items"]):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.bullet = True
        para.space_after = Pt(6)
        para.line_spacing = 1.1
        run = para.add_run()
        run.text = fill(item, 70)
        set_run_style(run, 18, False, TEXT)

    add_footer(slide, "Closing", "Operational checklist")


def write_speaker_script():
    order = {
        "Manufacturing Training Walk-through": 1,
        "Training Agenda": 2,
        "Home Dashboard": 3,
        "Manufacturing Overview": 4,
        "Product Master: FG-PSS-TH-01005": 5,
        "Product Inventory & Routes": 6,
        "Reordering Rule / Min-Max": 7,
        "Bill of Materials": 8,
        "Purchase Order from Shortage": 9,
        "Promotion / FOC Sales Order": 10,
        "SO -> MTO Example": 11,
        "Manufacturing Order from Replenishment": 12,
        "Work Order / Shop Floor": 13,
        "Workcenter Master: Injection 5": 14,
        "Workcenter Compatibility Matrix": 15,
        "Mold Master": 16,
        "Mold Matrix": 17,
        "Transfer Plastic": 18,
        "Transfer Pharma": 19,
        "Invoice & Accounting Endpoint": 20,
        "Closing Checklist": 21,
    }
    sorted_slides = sorted(SLIDES, key=lambda item: order.get(item["title"], 999))
    lines = []
    for index, data in enumerate(sorted_slides, start=1):
        title = data["title"]
        lines.append(f"Slide {index}: {title}")
        lines.append("=" * (len(lines[-1])))
        if data.get("menu"):
            lines.append(f"Menu path: {data['menu']}")
        if data.get("section"):
            lines.append(f"Section: {data['section']}")
        if data.get("what"):
            lines.append("What this page is for:")
            for item in data["what"]:
                lines.append(f"- {item}")
        if data.get("fields"):
            lines.append("Fields to explain:")
            for item in data["fields"]:
                lines.append(f"- {item}")
        if data.get("items") and data.get("kind") in {"agenda", "closing"}:
            lines.append("Items:")
            for item in data["items"]:
                lines.append(f"- {item}")
        if data.get("speaker"):
            lines.append("Speaker notes:")
            for item in data["speaker"]:
                lines.append(f"- {item}")
        lines.append("")

    OUT_TXT.write_text("\n".join(lines), encoding="utf-8")


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    order = {
        "Manufacturing Training Walk-through": 1,
        "Training Agenda": 2,
        "Home Dashboard": 3,
        "Manufacturing Overview": 4,
        "Product Master: FG-PSS-TH-01005": 5,
        "Product Inventory & Routes": 6,
        "Reordering Rule / Min-Max": 7,
        "Bill of Materials": 8,
        "Purchase Order from Shortage": 9,
        "Promotion / FOC Sales Order": 10,
        "SO -> MTO Example": 11,
        "Manufacturing Order from Replenishment": 12,
        "Work Order / Shop Floor": 13,
        "Workcenter Master: Injection 5": 14,
        "Workcenter Compatibility Matrix": 15,
        "Mold Master": 16,
        "Mold Matrix": 17,
        "Transfer Plastic": 18,
        "Transfer Pharma": 19,
        "Invoice & Accounting Endpoint": 20,
        "Closing Checklist": 21,
    }
    sorted_slides = sorted(SLIDES, key=lambda item: order.get(item["title"], 999))

    for data in sorted_slides:
        kind = data.get("kind", "content")
        if kind == "title":
            build_title(prs, data)
        elif kind == "agenda":
            build_agenda(prs, data)
        elif kind == "closing":
            build_closing(prs, data)
        else:
            build_content(prs, data)

    prs.save(OUT_PPTX)
    write_speaker_script()


if __name__ == "__main__":
    import atexit

    atexit.register(main)

SLIDES.extend(
    [
        {
            "title": "Promotion / FOC Sales Order",
            "section": "Flow 1",
            "menu": "Sales > Orders",
            "screenshot": "06_sale_foc_sob_263069.png",
            "what": [
                "ใช้สอน flow ส่งเสริมการขายและ free item ในเอกสารขายจริง",
                "ตัวอย่างนี้มี line ฟรีอยู่ใน order lines และยังมี flow downstream ต่อได้",
            ],
            "fields": [
                "`Customer`, `Invoice Address`, `Delivery Address` ใช้กำหนดปลายทางเอกสาร",
                "`SO Type` ใช้แยกประเภทคำสั่งขาย",
                "`Order Date` และ `Delivery Date` ใช้ประกอบการวางแผน fulfillment",
                "`Delivery` smart button ใช้ trace เอกสารส่งของต่อจาก SO",
                "`Order Lines` ต้องสังเกต free line ที่มีราคา 0 และ tag [FREE]",
                "`Untaxed Amount`, `VAT`, `Total` ใช้ยืนยันมูลค่าทางการขาย",
            ],
            "speaker": [
                "สอนให้ผู้เรียนเห็นว่า free line อยู่ใน SO อย่างไร ไม่ได้เป็นเอกสารพิเศษแยกออกไป",
                "ย้ำว่าถึงเป็น FOC ก็ยังมี stock movement และอาจมีผลด้าน accounting/costing",
                "ให้ผู้เรียนมอง smart button ด้านบนเป็นตัวเชื่อมไป Delivery หรือเอกสาร downstream",
            ],
        },
        {
            "title": "SO -> MTO Example",
            "section": "Flow 2",
            "menu": "Sales > Orders",
            "screenshot": "07_sale_mto_sob_263070.png",
            "what": [
                "ใช้สอน demand ที่เริ่มจาก sale แล้ววิ่งต่อไป production",
                "เทียบกับ FOC เพื่อให้เห็นว่าหน้าจอ SO คล้ายกัน แต่ intent ของ flow ต่างกัน",
            ],
            "fields": [
                "`Delivery` smart button ใช้ดู picking ที่ถูกสร้างจาก order นี้",
                "`Create Invoice` ใช้เรียก invoice ในกรณีที่ workflow กำหนดไว้",
                "`Order Lines` ใช้ดูสินค้า, quantity, delivered และ invoiced",
                "`Unit Price` ของ line ปกติเทียบกับ free line ช่วยสอนความต่างของ promotion",
                "`Statusbar` ใช้สอนว่าเอกสารผ่าน Quotation, To Approve, Sales Order อย่างไร",
            ],
            "speaker": [
                "ชี้ว่าตัวอย่างนี้ผู้ใช้งานไม่ต้องเปิด MO ลูกเอง ถ้า setting และ route ถูก",
                "ให้ผู้เรียนสังเกต line ปกติเทียบกับ line ฟรีและใช้จุดนี้อธิบายว่า promotion สามารถอยู่ร่วมกับ flow MTO ได้",
                "เชื่อมต่อไปสไลด์ MO เพื่อให้เห็น demand เดียวกันจากฝั่ง production",
            ],
        },
        {
            "title": "Manufacturing Order from Replenishment",
            "section": "Flow 3",
            "menu": "Manufacturing > Operations > Manufacturing Orders",
            "screenshot": "08_mo_mts_gmp_moph_00011.png",
            "what": [
                "ใช้สอน MO ที่มาจาก replenishment / Min-Max ไม่ได้มาจาก sale โดยตรง",
                "เป็นหน้าหลักที่ planner และ production ใช้ตรวจ component, workorders และต้นทุน",
            ],
            "fields": [
                "`Product` และ `Quantity To Produce` คือสิ่งที่จะผลิต",
                "`Bill of Material` ชี้กลับไปที่ master data ที่ใช้ในงานนี้",
                "`Scheduled Date` และ `Scheduled End` ใช้วางแผนเวลา",
                "`Component Status` ช่วยบอกว่า raw materials พร้อมหรือไม่",
                "`Transfers`, `Shop Floor`, `MO Cost`, `Workorders` smart buttons ใช้ trace flow ต่อ",
                "`Components` tab ใช้ดูรายการที่ต้อง consume พร้อม source location และ quantity to consume",
            ],
            "speaker": [
                "ชี้คำว่า generated via replenishment ใน chatter ถ้ามี เพื่อย้ำว่ามาจาก Min/Max",
                "สอนให้ดู Component Status ก่อนเสมอ เพราะถ้าวัตถุดิบไม่พร้อม จะเปิดงานจริงไม่ได้",
                "ย้ำว่า MO เป็นจุดรวมข้อมูลหลักของงาน ทั้ง component, workorder, transfer และ cost",
            ],
        },
        {
            "title": "Work Order / Shop Floor",
            "section": "Execution",
            "menu": "Manufacturing > Operations > Work Orders",
            "screenshot": "15_workorder_operation_112.png",
            "what": [
                "ใช้สอนหน้าปฏิบัติการจริงที่ operator หรือ supervisor เจอ",
                "เป็นจุดที่เริ่มทำงาน, บันทึกเวลา และดูสถานะพร้อมของ/components",
            ],
            "fields": [
                "`Work Center` บอกว่า operation นี้จะทำที่เครื่องไหน",
                "`Product` และ `Quantity` บอกว่ากำลังผลิตอะไรและปริมาณเท่าไร",
                "`Molds` และ `Mold Cost` สำคัญกับงานฝั่ง plastic หรือกรณีที่มี mold mapping",
                "`Manufacturing Order` ใช้ย้อนกลับไปเอกสารแม่",
                "`Expected Duration` คือเวลามาตรฐานที่ระบบคำนวณให้",
                "`Time Tracking` ใช้ดู employee productivity, start/end และเวลาจริง",
            ],
            "speaker": [
                "อธิบายว่าหน้านี้คือ operational truth ของหน้างาน ไม่ใช่แค่ข้อมูล master",
                "ถ้า operator ถามว่าเริ่มงานตรงไหน ให้ชี้ปุ่ม Open Shop Floor หรือ Mark as Done ตามสิทธิ์และ workflow",
                "สอนให้หัวหน้างานดู state ของ work order และ component readiness ก่อนเริ่มงาน",
            ],
        },
        {
            "title": "Workcenter Master: Injection 5",
            "section": "Workcenter",
            "menu": "Manufacturing > Configuration > Work Centers",
            "screenshot": "10_workcenter_injection5.png",
            "what": [
                "ใช้สอน master ของเครื่องจักรจริงฝั่ง Plastic",
                "หน้าจอนี้เชื่อมกับการวางแผน, costing และ allowed operations",
            ],
            "fields": [
                "`Work Center Name` และ `Code` ใช้อ้างอิงเครื่อง",
                "`Is Mold?` แยก machine ออกจาก mold",
                "`Manufacturing Type` ช่วยจัดกลุ่มเครื่องตามโรงงาน",
                "`Time Efficiency`, `Capacity`, `OEE Target` ใช้ในการวางแผนและวัดผล",
                "`Cost per hour` และ `per employee` ใช้เป็นฐานของ machine/labor cost",
                "`Analytic Distribution` ใช้โยงต้นทุนเข้า analytic/accounting",
            ],
            "speaker": [
                "ย้ำว่าหน้า workcenter ฝั่งเครื่องเป็น master data ของ machine ไม่ใช่ mold",
                "ชี้ว่าค่า cost per hour และ per employee มีผลโดยตรงกับต้นทุนจริงของ MO",
                "สอนให้ทีมแยกหน้าจอ machine กับ mold ให้ชัด เพราะ model เดียวกันแต่บทบาทต่างกัน",
            ],
        },
        {
            "title": "Workcenter Compatibility Matrix",
            "section": "Workcenter",
            "menu": "Manufacturing > Configuration > Work Centers > Compatibility Matrix",
            "screenshot": "10b_workcenter_injection5_compatibility_matrix.png",
            "what": [
                "ใช้สอนว่าเครื่องหนึ่งเครื่องรับ mold อะไรได้บ้าง",
                "เป็นแกนของ auto machine/mold selection ฝั่ง plastic",
            ],
            "fields": [
                "`Work Center` ในตารางคือ mold ที่เครื่องนี้ใช้ร่วมกันได้",
                "`Mold Status` บอกสถานะของ mold เช่น Normal, Warning, Full",
                "`Usage` แสดงจำนวน shots ที่ใช้ไปแล้ว",
                "`Limit` คือ mold life limit ที่ใช้เตือนการบำรุงรักษา",
                "`Add a line` ใช้เพิ่ม mold compatibility ใหม่ให้เครื่อง",
            ],
            "speaker": [
                "อธิบายว่าถ้าไม่มี record ใน compatibility matrix ระบบจะ auto assign mold ให้เครื่องนี้ไม่ได้",
                "ชี้การอ่าน Usage เทียบกับ Limit ว่าหัวหน้างานควรใช้เฝ้าความเสี่ยงเรื่อง mold life อย่างไร",
                "สไลด์นี้ใช้เชื่อมต่อไปหน้าฝั่ง mold ซึ่งเป็นอีกด้านของ matrix เดียวกัน",
            ],
        },
    ]
)
