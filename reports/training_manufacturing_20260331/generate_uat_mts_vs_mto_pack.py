from __future__ import annotations

from pathlib import Path
from textwrap import dedent
from xml.sax.saxutils import escape

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(r"C:\365_project\TheCool18e\Dev\reports\training_manufacturing_20260331")

PPTX_OUT = BASE / "manufacturing_uat_mts_vs_mto_20260403.pptx"
FLOW_MD_OUT = BASE / "11_mts_vs_mto_flow_diagram_th.md"
FLOW_SVG_OUT = BASE / "11_mts_vs_mto_flow_diagram_th.svg"
TRAINER_OUT = BASE / "12_mts_vs_mto_trainer_script_th.md"


SLIDES = [
    {
        "title": "MTS vs MTO ใน UAT",
        "subtitle": "เปรียบเทียบตามลำดับ Replenishment -> Route -> Rule -> Operation Type -> Putaway",
        "bullets": [
            "ตัวอย่าง MTS: FG-PNC-TH-01001",
            "ตัวอย่าง MTO: FG-MTK-IL-01001",
            "ใช้เพื่ออธิบายว่าความต่างจริงอยู่ที่จุดเริ่ม demand และ chain ที่ระบบสร้าง",
        ],
    },
    {
        "title": "เลือกตัวอย่างจาก UAT อย่างไร",
        "bullets": [
            "FG-PNC-TH-01001 มี orderpoint 175 ที่ GMP/Stock, trigger auto, ไม่มี MTO route",
            "FG-MTK-IL-01001 มี route Replenish on Order (MTO) และไม่มี orderpoint",
            "สองตัวนี้จึงแยกกันชัด: ตัวแรกเติม stock, ตัวหลังตอบสนองต่อ sales demand",
        ],
    },
    {
        "title": "MTS: FG-PNC-TH-01001",
        "bullets": [
            "Replenishment = orderpoint 175 ที่ GMP/Stock, min 0, max 0, trigger auto",
            "Route = Manufacture (Pharma)",
            "Rule = 146, action manufacture, source GMP/Stock/คลังลอย, destination GMP/Stock",
            "Operation Type = Manufacturing Pharma",
            "Putaway = ใช้กับของที่เข้าคลังระหว่างทางหรือ component ที่รับเข้า ไม่ใช่ตัว trigger",
        ],
    },
    {
        "title": "MTO: FG-MTK-IL-01001",
        "bullets": [
            "Replenishment = sales demand จาก SO ไม่ใช่ orderpoint",
            "Route = Replenish on Order (MTO) + Manufacture (Pharma) + Auto Transfer Semi (Pharma)",
            "Rule = 5 สำหรับ pull จากลูกค้า, 146 สำหรับผลิต, 144 สำหรับโอน semi",
            "Operation Type = Pick, Transfer Pharma, Manufacturing Pharma",
            "Putaway = มีผลกับ bought components และ semi ที่เข้าคลังระหว่างทาง ไม่ใช่ขา customer",
        ],
    },
    {
        "title": "ความต่างที่ต้องสอนทีมให้เห็น",
        "bullets": [
            "MTS เริ่มจาก stock shortage ในคลัง",
            "MTO เริ่มจาก sales demand ของลูกค้า",
            "MTS มองปลายทางเป็นการเติมกลับเข้า stock",
            "MTO มองปลายทางเป็นการตอบออเดอร์ และ supply chain จะถูกดึงย้อนกลับ",
        ],
    },
    {
        "title": "Rule และ Operation Type ที่เกี่ยวข้อง",
        "bullets": [
            "Rule 146 -> Manufacturing Pharma",
            "Rule 145 -> Transfer Plastic",
            "Rule 144 -> Transfer Pharma",
            "Rule 7 -> Buy / Receipts",
            "Rule 5 -> Pick ฝั่ง MTO จาก GMP/Stock ไป Customers",
        ],
    },
    {
        "title": "สรุปแบบภาษาหน้างาน",
        "bullets": [
            "MTS = ของในคลังไม่พอ ระบบเติม stock ให้",
            "MTO = ลูกค้าสั่งก่อน ระบบค่อยย้อนกลับไปหาของให้",
            "อย่าดูแค่ชื่อ route ต้องดูว่า demand เริ่มจากไหนด้วย",
            "เวลาระบบไม่ไปต่อ ให้ไล่ Trigger -> Route -> Rule -> Operation Type -> Putaway",
        ],
    },
]


FLOW_MD = dedent(
    """
    # Flow Diagram 2 คอลัมน์: MTS vs MTO ใน UAT

    ## ตัวอย่างที่ใช้
    - MTS: `FG-PNC-TH-01001`
    - MTO: `FG-MTK-IL-01001`

    ## สรุปสั้น
    - `MTS` = ของในคลังขาดก่อน แล้ว orderpoint ยิง procurement ไปเติม stock
    - `MTO` = ลูกค้าสั่งก่อน แล้ว demand ดึง supply chain ย้อนกลับไปสร้างของ

    ## MTS Column
    1. `Replenishment`
       - Trigger จาก orderpoint `175`
       - Location = `GMP/Stock`
       - min = `0`, max = `0`, trigger = `auto`
    2. `Route`
       - Product route = `Manufacture (Pharma)`
    3. `Rule`
       - Rule `146`
       - action = `manufacture`
       - source = `GMP/Stock/คลังลอย`
       - destination = `GMP/Stock`
    4. `Operation Type`
       - `Manufacturing Pharma`
    5. `Putaway`
       - ทำงานตอน stock/compoent เข้าปลายทางใหญ่แล้ว
       - ไม่ได้เป็นตัว trigger procurement

    ## MTO Column
    1. `Replenishment`
       - Trigger จาก `SO demand`
       - ไม่มี orderpoint
    2. `Route`
       - `Replenish on Order (MTO)`
       - `Manufacture (Pharma)`
       - `Auto Transfer Semi (Pharma)`
    3. `Rule`
       - Rule `5` = pull จาก `GMP/Stock` ไป `Customers`
       - Rule `146` = manufacture ฝั่ง pharma
       - Rule `144` = transfer semi pharma
    4. `Operation Type`
       - `Pick`
       - `Transfer Pharma`
       - `Manufacturing Pharma`
    5. `Putaway`
       - ไปมีผลกับ material/semi ที่เข้าคลังระหว่างทาง
       - ไม่ได้อยู่บนขา customer โดยตรง

    ## สิ่งที่ต้องเน้นเวลาเทียบ
    - MTS เริ่มจาก `stock shortage`
    - MTO เริ่มจาก `sales demand`
    - MTS เป้าหมายคือ `เติม stock`
    - MTO เป้าหมายคือ `ตอบออเดอร์`

    ## Mermaid
    ```mermaid
    flowchart LR
        subgraph MTS["MTS : FG-PNC-TH-01001"]
            A1[Replenishment\\nOrderpoint 175\\nGMP/Stock] --> A2[Route\\nManufacture (Pharma)]
            A2 --> A3[Rule 146\\naction = manufacture]
            A3 --> A4[Operation Type\\nManufacturing Pharma]
            A4 --> A5[Putaway / Storage\\nทำงานตอนของเข้าคลัง]
        end

        subgraph MTO["MTO : FG-MTK-IL-01001"]
            B1[Replenishment\\nSO Demand\\nไม่มี orderpoint] --> B2[Route\\nMTO + Manufacture (Pharma) + Auto Transfer Semi (Pharma)]
            B2 --> B3[Rule 5 / 146 / 144\\npull + manufacture + transfer]
            B3 --> B4[Operation Type\\nPick + Transfer Pharma + Manufacturing Pharma]
            B4 --> B5[Putaway / Storage\\nใช้กับ material ระหว่างทาง]
        end
    ```
    """
).strip() + "\n"


TRAINER_MD = dedent(
    """
    # Trainer Script ภาษาไทย: MTS vs MTO ใน UAT

    ## เปิดเรื่อง
    วันนี้เราจะเทียบ MTS กับ MTO แบบเห็นเป็นรูปธรรม โดยใช้สินค้าจริงจาก UAT คนละตัว เพื่อให้ทีมเห็นว่าความต่างไม่ได้อยู่แค่ชื่อ route แต่ต่างกันตั้งแต่ตัวจุดชนวนของ demand จนถึงเอกสารที่ระบบสร้าง

    ตัวอย่างที่ใช้:
    - MTS = FG-PNC-TH-01001
    - MTO = FG-MTK-IL-01001

    ---

    ## Slide 1: หัวข้อ
    ให้บอกผู้เรียนก่อนว่าเราจะใช้กรอบเดียวกันทั้งสองฝั่ง:
    Replenishment -> Route -> Rule -> Operation Type -> Putaway

    จุดประสงค์คือให้ตอบได้ว่า:
    - ทำไมสินค้าตัวหนึ่งเติม stock
    - ทำไมอีกตัวหนึ่งรอ demand จาก sales

    ## Slide 2: ทำไมเลือกสองตัวนี้
    อธิบายว่า:
    - FG-PNC-TH-01001 มี orderpoint 175 จริง และไม่มี MTO route
    - FG-MTK-IL-01001 มี MTO route จริง และไม่มี orderpoint

    เพราะฉะนั้นมันเป็นตัวอย่างสะอาดสำหรับอธิบายความต่าง

    ## Slide 3: MTS
    พูดตามลำดับนี้
    1. ตัวเริ่มคือ orderpoint 175 ที่ GMP/Stock
    2. เมื่อ forecast ขาด ระบบเริ่ม procurement อัตโนมัติ
    3. Product route ของมันคือ Manufacture (Pharma)
    4. ดังนั้น rule ที่ทำงานจริงคือ rule 146
    5. เอกสารที่คนเห็นคือ Manufacturing Pharma

    ประโยคที่ควรย้ำ:
    MTS คิดจาก stock ก่อน ถ้าของในคลังไม่พอ ระบบจึงค่อยสร้างงานมาเติม

    ## Slide 4: MTO
    พูดตามลำดับนี้
    1. ตัวเริ่มไม่ใช่ orderpoint แต่เป็น demand จาก SO
    2. Product มี MTO route อยู่จริง
    3. MTO route ดึง demand ย้อนกลับผ่าน rule 5
    4. ถ้าต้องผลิต ระบบจะต่อไป rule 146
    5. ถ้าต้องดึง semi ก็จะต่อไป rule 144

    ประโยคที่ควรย้ำ:
    MTO ไม่ได้เติม stock ล่วงหน้า แต่ค่อยสร้าง supply เมื่อมี order เข้ามา

    ## Slide 5: ความต่างที่ต้องให้ทีมจำ
    ให้สรุปด้วยคำง่าย ๆ
    - MTS = stock-first
    - MTO = order-first
    - MTS เป้าหมายคือเติม stock
    - MTO เป้าหมายคือส่งมอบตามออเดอร์

    ## Slide 6: Rule และ Operation Type
    ใช้สไลด์นี้ย้ำว่า route อย่างเดียวไม่พอ
    ต้องดู rule และ operation type ด้วย

    ให้พูดว่า:
    - ถ้าเป็น MTS ตัวอย่างนี้ คนหน้างานจะเห็น Manufacturing Pharma เป็นหลัก
    - ถ้าเป็น MTO คนหน้างานจะเห็น Pick, Transfer Pharma, Manufacturing Pharma ตามชั้นของ chain

    ## Slide 7: สรุปปิด
    ปิดด้วยประโยคนี้:
    ถ้าถามว่าของขาดแล้วระบบจะทำอะไร อย่าตอบจากความเคยชิน ให้ไล่ถาม 5 เรื่องนี้เสมอ คือ trigger มาจากไหน, route อะไร, rule ไหนทำงาน, operation type อะไรจะถูกสร้าง, และ putaway จะพาไปเก็บที่ไหน

    ---

    ## Q&A ที่คนเรียนมักถาม
    ### ถาม: สินค้ามีทั้ง orderpoint และ MTO route ได้ไหม
    ตอบ: ได้ แต่จะทำให้ logic ซ้อนกัน ต้องดูว่า demand จริงเริ่มจาก orderpoint หรือ sales flow และต้องดู route ที่ level product กับ orderpoint พร้อมกัน

    ### ถาม: min/max = 0 แปลว่าไม่เติมของใช่ไหม
    ตอบ: ไม่ใช่ ในฐานนี้ 0/0 ถูกใช้แบบเติม shortage กลับมาที่ 0

    ### ถาม: MTO แล้วถ้ามีของใน stock อยู่จะใช้ stock ได้ไหม
    ตอบ: ขึ้นกับ rule และ procure_method ที่วิ่งจริง บาง rule ใน UAT เป็น `mts_else_mto` จึงยังมีพฤติกรรมใช้ stock ก่อนในบาง leg

    ### ถาม: Putaway มีผลกับ MTO ไหม
    ตอบ: มีในส่วนของ material หรือ semi ที่เข้าคลังระหว่างทาง แต่ไม่ได้เป็นตัวเริ่ม MTO demand
    """
).strip() + "\n"


def set_run_font(run, size, bold=False, color=None):
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_title(slide, title, subtitle=None):
    tb = slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12), Inches(1.1))
    tf = tb.text_frame
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    set_run_font(r, 24, bold=True, color=RGBColor(27, 47, 77))
    if subtitle:
        p = tf.add_paragraph()
        r = p.add_run()
        r.text = subtitle
        set_run_font(r, 11, color=RGBColor(92, 102, 115))


def add_bullets(slide, bullets):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.7), Inches(1.55), Inches(11.85), Inches(5.35)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(245, 248, 252)
    shape.line.color.rgb = RGBColor(208, 219, 231)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.space_after = Pt(9)
        r = p.add_run()
        r.text = bullet
        set_run_font(r, 18, color=RGBColor(38, 50, 56))


def build_ppt():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    for slide_spec in SLIDES:
        slide = prs.slides.add_slide(prs.slide_layouts[6])
        add_title(slide, slide_spec["title"], slide_spec.get("subtitle"))
        add_bullets(slide, slide_spec["bullets"])
    prs.save(PPTX_OUT)


def build_svg():
    parts = [
        '<svg xmlns="http://www.w3.org/2000/svg" width="1200" height="760" viewBox="0 0 1200 760">',
        '<defs><marker id="arrow" markerWidth="10" markerHeight="10" refX="8" refY="3" orient="auto" markerUnits="strokeWidth"><path d="M0,0 L0,6 L9,3 z" fill="#6E7F99"/></marker></defs>',
        '<rect width="100%" height="100%" fill="#ffffff"/>',
        '<text x="40" y="40" font-family="Tahoma" font-size="26" font-weight="700" fill="#1B2F4D">Flow Diagram 2 คอลัมน์: MTS vs MTO (UAT)</text>',
        '<text x="120" y="95" font-family="Tahoma" font-size="22" font-weight="700" fill="#355C7D">MTS : FG-PNC-TH-01001</text>',
        '<text x="680" y="95" font-family="Tahoma" font-size="22" font-weight="700" fill="#355C7D">MTO : FG-MTK-IL-01001</text>',
    ]

    left_boxes = [
        (80, 120, 430, 90, "#E8F1FF", "Replenishment\nOrderpoint 175\nGMP/Stock, auto, min/max = 0/0"),
        (80, 250, 430, 90, "#EEF7EA", "Route\nManufacture (Pharma)"),
        (80, 380, 430, 110, "#FFF5E6", "Rule\n146 = manufacture\nsource GMP/Stock/คลังลอย\ndestination GMP/Stock"),
        (80, 530, 430, 90, "#F6EEFF", "Operation Type\nManufacturing Pharma"),
        (80, 660, 430, 60, "#F7F7F7", "Putaway\nทำงานตอน stock / component เข้าปลายทาง"),
    ]

    right_boxes = [
        (640, 120, 480, 90, "#E8F1FF", "Replenishment\nSO Demand\nไม่มี orderpoint"),
        (640, 250, 480, 90, "#EEF7EA", "Route\nMTO + Manufacture (Pharma) + Auto Transfer Semi (Pharma)"),
        (640, 380, 480, 110, "#FFF5E6", "Rule\n5 = pull จาก GMP/Stock ไป Customers\n146 = manufacture\n144 = transfer semi"),
        (640, 530, 480, 90, "#F6EEFF", "Operation Type\nPick + Transfer Pharma + Manufacturing Pharma"),
        (640, 660, 480, 60, "#F7F7F7", "Putaway\nใช้กับ material / semi ระหว่างทาง"),
    ]

    for x, y, w, h, fill, text in left_boxes + right_boxes:
        parts.append(
            f'<rect x="{x}" y="{y}" width="{w}" height="{h}" rx="14" ry="14" fill="{fill}" stroke="#B8C7D9" stroke-width="1.5"/>'
        )
        for i, line in enumerate(text.split("\n")):
            parts.append(
                f'<text x="{x+16}" y="{y+28+i*20}" font-family="Tahoma" font-size="17" fill="#263238">{escape(line)}</text>'
            )

    for x in (295, 880):
        parts.append(f'<line x1="{x}" y1="210" x2="{x}" y2="250" stroke="#6E7F99" stroke-width="2" marker-end="url(#arrow)"/>')
        parts.append(f'<line x1="{x}" y1="340" x2="{x}" y2="380" stroke="#6E7F99" stroke-width="2" marker-end="url(#arrow)"/>')
        parts.append(f'<line x1="{x}" y1="490" x2="{x}" y2="530" stroke="#6E7F99" stroke-width="2" marker-end="url(#arrow)"/>')
        parts.append(f'<line x1="{x}" y1="620" x2="{x}" y2="660" stroke="#6E7F99" stroke-width="2" marker-end="url(#arrow)"/>')

    parts.append('<text x="80" y="742" font-family="Tahoma" font-size="15" fill="#4B5968">MTS = stock-first, เติม stock เมื่อขาด</text>')
    parts.append('<text x="640" y="742" font-family="Tahoma" font-size="15" fill="#4B5968">MTO = order-first, ลูกค้าสั่งก่อนแล้วค่อยย้อนกลับไปสร้าง supply</text>')
    parts.append("</svg>")
    FLOW_SVG_OUT.write_text("\n".join(parts), encoding="utf-8")


def main():
    FLOW_MD_OUT.write_text(FLOW_MD, encoding="utf-8")
    TRAINER_OUT.write_text(TRAINER_MD, encoding="utf-8")
    build_svg()
    build_ppt()
    print(FLOW_MD_OUT)
    print(FLOW_SVG_OUT)
    print(TRAINER_OUT)
    print(PPTX_OUT)


if __name__ == "__main__":
    main()
