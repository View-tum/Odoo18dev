from pathlib import Path
from textwrap import fill

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


BASE = Path(r"C:\365_project\TheCool18e\Dev\reports\training_manufacturing_20260331")
SHOT_DIR = BASE / "uat_screenshots"
ANNOTATED_DIR = BASE / "uat_screenshots_annotated"
OUT_PPTX = BASE / "manufacturing_training_uat_detailed_20260331.pptx"
OUT_TXT = BASE / "manufacturing_training_uat_detailed_speaker_20260331.txt"

BG = RGBColor(244, 247, 249)
WHITE = RGBColor(255, 255, 255)
GREEN = RGBColor(21, 101, 52)
GREEN_DARK = RGBColor(16, 72, 40)
TEXT = RGBColor(32, 41, 56)
MUTED = RGBColor(89, 98, 111)
LINE = RGBColor(208, 214, 221)


def clean_text(value):
    if not isinstance(value, str):
        return value
    if any(marker in value for marker in ("à", "â", "Ã")):
        for encoding in ("latin1", "cp1252"):
            try:
                return value.encode(encoding).decode("utf-8")
            except (UnicodeEncodeError, UnicodeDecodeError):
                continue
    return value


SLIDES = [
    {
        "kind": "title",
        "title": "Manufacturing Training Walk-through",
        "subtitle": "UAT Database\nDetailed click-by-click training\nPlastic + Pharma + PO + Shop Floor + Mold + Costing",
        "speaker": [
            "เปิดคลาสด้วยการย้ำว่าชุดนี้ใช้ภาพจริงจาก UAT และมีเลขกำกับจุดที่ต้องกดทุกหน้า",
            "บอกผู้เรียนว่าวันนี้จะไม่สอนแค่แนวคิด แต่จะพาเดินตามหน้าจอจริงตั้งแต่ต้นจนปลาย",
        ],
    },
    {
        "kind": "agenda",
        "title": "Training Agenda",
        "items": [
            "1. ภาพรวมระบบและแผนที่เอกสาร",
            "2. Product, Routes, Min-Max และ BoM",
            "3. Purchasing จาก shortage",
            "4. Promotion / SO 0 บาท",
            "5. SO -> MTO และ MO จาก Replenishment",
            "6. GMP Shop Floor, Work Order และ Scrap",
            "7. Workcenter และ Mold Mapping",
            "8. Transfer Plastic / Pharma",
            "9. Invoice, Costing และ Closing Checklist",
        ],
        "speaker": [
            "ใช้สไลด์นี้บอก flow ของคลาสว่าเริ่มจาก master data แล้วไปจบที่ transaction ปลายทาง",
            "ย้ำว่าทุกหัวข้อจะมีภาพจริงพร้อมตำแหน่งที่ต้องกด ไม่ใช่สรุปเชิงทฤษฎีอย่างเดียว",
        ],
    },
]

SLIDES.extend(
    [
        {
            "kind": "content",
            "title": "Home Dashboard",
            "section": "Opening",
            "menu": "Home",
            "image": "01_home_dashboard.png",
            "purpose": [
                "ใช้หน้า Home เป็นจุดตั้งต้นให้ผู้เรียนเห็นว่า flow งานผลิตเชื่อมหลายแอป",
                "สอนว่าถ้าต้อง trace ปัญหา ต้องรู้ว่าแต่ละทีมเข้าเมนูไหน",
            ],
            "clicks": [
                "1. เริ่มจากไอคอน Manufacturing เพื่อเข้าเมนูหลักของงานผลิต",
                "2. ใช้ Inventory เมื่อต้องตรวจ stock, transfer, scrap และ reordering rule",
                "3. ใช้ Purchase เมื่อต้องตาม RFQ หรือ PO จากของที่ route เป็น Buy",
                "4. ใช้ GMP Shop Floor เมื่อต้องดูงานหน้างานแบบ card และเปิด console",
                "5. ใช้ Accounting เมื่อต้องตาม invoice และผลทางบัญชีปลายทาง",
            ],
            "fields": [
                "Manufacturing คือศูนย์กลางของ MO, WO, BoM และ workcenter",
                "Inventory คือที่ดู availability, movements, locations และ min/max",
                "Purchase คือที่ดูผลของ shortage ฝั่ง Buy route",
                "GMP Shop Floor คือหน้าปฏิบัติการจริงสำหรับหัวหน้างานและ operator",
                "Accounting คือปลายทางของ sale invoice และผลกระทบด้านบัญชี",
            ],
            "annotations": [
                {"n": 1, "label": (520, 365), "target": (520, 365)},
                {"n": 2, "label": (195, 365), "target": (195, 365)},
                {"n": 3, "label": (1040, 365), "target": (1040, 365)},
                {"n": 4, "label": (365, 365), "target": (365, 365)},
                {"n": 5, "label": (860, 365), "target": (860, 365)},
            ],
            "speaker": [
                "อธิบายว่าหน้านี้เป็นแผนที่ระบบ ไม่ใช่แค่หน้ารวมแอป",
                "ให้ผู้เข้าอบรมจำความสัมพันธ์ Manufacturing, Inventory, Purchase, Shop Floor และ Accounting ให้ได้ก่อนเข้า flow รายละเอียด",
            ],
        },
        {
            "kind": "content",
            "title": "Manufacturing Overview",
            "section": "Opening",
            "menu": "Manufacturing > Overview",
            "image": "02_manufacturing_overview.png",
            "purpose": [
                "อธิบายโครงเมนู Manufacturing และหน้าภาพรวมที่ planner กับ supervisor ใช้เปิดทุกวัน",
                "สอนให้ผู้เรียนเห็นความต่างระหว่างงานปฏิบัติการ, planning และ configuration",
            ],
            "clicks": [
                "1. เข้า Overview เพื่อดูภาพรวม work center และภาระงาน",
                "2. ใช้เมนู Operations เมื่อต้องเจาะไป MO หรือ WO",
                "3. ใช้เมนู Products เมื่อต้องย้อนกลับไปดู product หรือ BoM",
                "4. ใช้เมนู Configuration เมื่อต้องแก้ workcenter, operation type หรือ master data",
            ],
            "fields": [
                "Overview ใช้ดูภาระงานของเครื่องและสถานะงานแบบภาพรวม",
                "Operations รวม MO, WO และ execution documents",
                "Products ใช้ดู product master, BoM และ related data",
                "Configuration ใช้ตั้งค่า workcenter, operation type และ logic ต้นทางของระบบ",
            ],
            "annotations": [
                {"n": 1, "label": (180, 35), "target": (180, 35)},
                {"n": 2, "label": (240, 35), "target": (240, 35)},
                {"n": 3, "label": (360, 35), "target": (360, 35)},
                {"n": 4, "label": (505, 35), "target": (505, 35)},
            ],
            "speaker": [
                "ย้ำว่านี่คือหน้า supervisor view ของ Manufacturing",
                "ถ้าจะไล่ปัญหางานผลิต ให้เริ่มจากภาพรวมแล้วค่อย drill ลง MO หรือ work center",
            ],
        },
        {
            "kind": "content",
            "title": "Product Master: FG-PSS-TH-01005",
            "section": "Master Data",
            "menu": "Manufacturing > Products > Products",
            "image": "03_product_fg_pss_th_01005.png",
            "purpose": [
                "ใช้ finished good ตัวจริงใน UAT เป็นตัวอย่างเพื่อสอนว่า product master คุม flow หลายเรื่องพร้อมกัน",
                "หน้าเดียวนี้เชื่อมไป inventory, bom, sale และ costing ได้",
            ],
            "clicks": [
                "1. เปิดสินค้าแล้วดูชื่อสินค้าและรหัสอ้างอิงให้ตรงกับของจริง",
                "2. เช็ก smart buttons ด้านบน เช่น On Hand, Forecasted และ Bill of Materials",
                "3. ดู Product Type และ Manufacturing Type เพื่อยืนยันว่าสินค้านี้เป็นฝั่ง Pharma",
                "4. ดู Sales Price และ Cost เพื่อเข้าใจ master data ฝั่งราคา",
                "5. ดู Category และ Reference เพื่อใช้ trace category-based logic และ report",
            ],
            "fields": [
                "Product Type ระบุว่าสินค้าเป็น storable goods, service หรือ combo",
                "Manufacturing Type ใช้แยก Plastic, Pharma หรือ Packaging และพา flow ไป operation type ที่ถูก",
                "Track Inventory บอกว่าติดตามเป็น lot, serial หรือไม่ติดตาม",
                "Invoicing Policy บอกว่าจะ invoice ตาม delivered หรือ ordered quantities",
                "Sales Price เป็นราคาขายมาตรฐาน ส่วน Cost เป็นต้นทุนอ้างอิง master data",
                "Category และ Reference มีผลกับ route inheritance, accounting และ reporting",
            ],
            "annotations": [
                {"n": 1, "label": (285, 160), "target": (285, 160)},
                {"n": 2, "label": (520, 52), "target": (520, 52)},
                {"n": 3, "label": (205, 346), "target": (205, 346)},
                {"n": 4, "label": (725, 345), "target": (725, 345)},
                {"n": 5, "label": (705, 492), "target": (705, 492)},
            ],
            "speaker": [
                "หน้า product ต้องสอนช้า เพราะเป็นต้นทางของการตัดสินใจทั้ง route, inventory และ accounting",
                "ย้ำว่า smart buttons ด้านบนใช้ตรวจของจริงใน stock และ BoM ได้ทันที ไม่ต้องออกจากหน้า",
            ],
        },
        {
            "kind": "content",
            "title": "Product Inventory & Routes",
            "section": "Master Data",
            "menu": "Manufacturing > Products > Products > Inventory",
            "image": "03b_product_fg_pss_th_01005_inventory_routes.png",
            "purpose": [
                "สอน route ซึ่งเป็นหัวใจของคำถามว่า ของขาดแล้วระบบจะไป MO หรือ PO",
                "ใช้หน้าจริงใน UAT เพื่อย้ำว่า FG ตัวนี้เป็น MTS/Manufacture ไม่ใช่ sale-driven MTO ตรง ๆ",
            ],
            "clicks": [
                "1. เข้าแท็บ Inventory ก่อนอธิบาย route",
                "2. ดูรายการ Routes ที่ถูกติ๊กจริงบนสินค้า",
                "3. ใช้ View Diagram เมื่อต้องอธิบาย logic procurement ให้ผู้ใช้เห็นภาพ",
                "4. ดู Category Routes ด้านล่างเพื่ออธิบาย route ที่ inherited มาจากระดับ category",
            ],
            "fields": [
                "Routes คือ route ที่ product ตัวนี้ใช้จริง เช่น Manufacture (Pharma), Buy หรือ Replenish on Order",
                "Category Routes คือ route ที่ inherited จาก category หรือ warehouse rules",
                "Responsible ใช้ระบุ owner ของข้อมูลหรือ logistics responsibility",
                "Lead Times ใช้ประกอบการวางแผนด้านเวลา ไม่ได้เป็นตัว trigger procurement โดยตรง",
                "Weight, Volume, Cartons และ DIMS สำคัญกับ logistics และ shipping planning",
            ],
            "annotations": [
                {"n": 1, "label": (245, 278), "target": (245, 278)},
                {"n": 2, "label": (171, 452), "target": (171, 452)},
                {"n": 3, "label": (162, 690), "target": (162, 690)},
                {"n": 4, "label": (202, 720), "target": (202, 720)},
            ],
            "speaker": [
                "ตรงนี้ต้องเน้นความต่างระหว่าง product routes กับ category routes ให้ชัด",
                "ให้ผู้เรียนตอบให้ได้ว่าทำไม FG-PSS-TH-01005 ถึงไป replenishment ไม่ได้สร้าง MO จาก SO แบบ MTO ตรง ๆ",
            ],
        },
        {
            "kind": "content",
            "title": "Reordering Rule / Min-Max",
            "section": "Planning",
            "menu": "Inventory > Configuration > Reordering Rules",
            "image": "04_reordering_rule_fg_pss_th_01005.png",
            "purpose": [
                "สอน MTS และ Min/Max จากข้อมูลจริงของ UAT",
                "ทำให้ทีมเข้าใจว่าถ้าไม่ได้เริ่มจาก SO ระบบจะสร้าง procurement จากจุดนี้อย่างไร",
            ],
            "clicks": [
                "1. เปิด Reordering Rule ของสินค้าและ location ที่ต้องการ",
                "2. ดู Min Quantity และ Max Quantity ว่าระบบจะเติม stock อย่างไร",
                "3. ดู Route และ Trigger เพื่อเข้าใจว่าจะไป Manufacture หรือ Buy",
                "4. ใช้ Forecast Description เพื่อตามว่าความต้องการมาจากเอกสารไหน",
            ],
            "fields": [
                "Min Quantity คือระดับต่ำสุดที่ยอมให้คงอยู่ใน location",
                "Max Quantity คือเป้าหมายที่ระบบจะเติมกลับไปหลังเห็น shortage",
                "Quantity Multiple ใช้บังคับให้เติมเป็นล็อตคูณ",
                "Route ของ orderpoint บอกว่าขาดแล้วจะไป MO หรือ PO",
                "Trigger และ Forecast Description ช่วย trace ว่าระบบยิง procurement เพราะอะไร",
            ],
            "annotations": [
                {"n": 1, "label": (220, 188), "target": (220, 188)},
                {"n": 2, "label": (545, 188), "target": (545, 188)},
                {"n": 3, "label": (830, 188), "target": (830, 188)},
                {"n": 4, "label": (1160, 188), "target": (1160, 188)},
            ],
            "speaker": [
                "ย้ำว่าบางกรณีใช้ 0/0 เพื่อเติม shortage กลับมาที่ศูนย์ได้ ไม่จำเป็นต้องมี safety stock สูงเสมอไป",
                "ให้ผู้เรียนแยกให้ออกว่า MTO กับ MTS ต่างกันที่ trigger ไม่ใช่แค่ชื่อ flow",
            ],
        },
        {
            "kind": "content",
            "title": "Bill of Materials",
            "section": "Master Data",
            "menu": "Manufacturing > Products > Bills of Materials",
            "image": "05_bom_fg_pss_th_01005.png",
            "purpose": [
                "สอน BoM ของ FG จริงใน UAT เพื่ออธิบายว่า MO ตัวแม่จะดึง component อะไรบ้าง",
                "ใช้เป็นจุดเชื่อมไป child MO, ซื้อ packaging และ route ของชิ้นส่วน",
            ],
            "clicks": [
                "1. เปิด BoM ของสินค้าให้ถูกตัว",
                "2. ดู Components ว่าสินค้าตัวแม่ต้องใช้ชิ้นส่วนอะไรบ้าง",
                "3. สลับไปดู Operations เมื่อต้องอธิบายขั้นตอนการผลิต",
                "4. เปิด Miscellaneous เมื่อต้องดู operation type หรือ field ประกอบอื่น",
            ],
            "fields": [
                "Product คือสินค้าที่ BoM นี้ใช้ผลิต",
                "Quantity คือจำนวนมาตรฐานต่อหนึ่งชุดผลิต",
                "BoM Type บอกว่าเป็น Manufacture this product, Kit หรือ Subcontracting",
                "Components table คือรายการวัตถุดิบหรือ semi ที่ต้องใช้จริง",
                "Operations คือขั้นตอนงานที่ไปสร้าง work orders",
            ],
            "annotations": [
                {"n": 1, "label": (160, 133), "target": (160, 133)},
                {"n": 2, "label": (165, 317), "target": (165, 317)},
                {"n": 3, "label": (160, 255), "target": (160, 255)},
                {"n": 4, "label": (230, 255), "target": (230, 255)},
            ],
            "speaker": [
                "ให้ผู้เรียนใช้ BoM เป็นจุดตั้งต้นเวลาอธิบายว่า manual MO จะดึงอะไรบ้าง",
                "ย้ำว่าถ้าชิ้นส่วนตัวใด route เป็น Buy ก็จะโยงไป PO, ถ้าเป็น Manufacture ก็จะไป child MO",
            ],
        },
    ]
)

SLIDES.extend(
    [
        {
            "kind": "content",
            "title": "Workcenter Master: Injection 5",
            "section": "Workcenter",
            "menu": "Manufacturing > Configuration > Work Centers",
            "image": "10_workcenter_injection5.png",
            "purpose": [
                "สอน master data ของเครื่องจริงฝั่ง plastic",
                "เป็นหน้าที่เชื่อม capacity, costing และ machine-level setup เข้าด้วยกัน",
            ],
            "clicks": [
                "1. เปิด workcenter ของเครื่องจริง",
                "2. ดู Is Mold? เพื่อแยกว่า record นี้เป็น machine ไม่ใช่ mold",
                "3. ดู Manufacturing Type, Capacity และ OEE Target",
                "4. ดู Cost per hour และ cost per employee เมื่อต้องสอนเรื่องต้นทุนเครื่อง",
                "5. ดูแท็บ Compatibility Matrix เพื่อเชื่อมไป mold mapping",
            ],
            "fields": [
                "Work Center Name และ Code ใช้อ้างอิงเครื่อง",
                "Is Mold? ใช้แยก machine ออกจาก mold แม้อยู่ model เดียวกัน",
                "Manufacturing Type ช่วยจัดกลุ่มเครื่องตามโรงงาน",
                "Capacity, Time Efficiency และ OEE Target ใช้วางแผนและวัดประสิทธิภาพ",
                "Cost per hour และ per employee มีผลกับ machine/labor costing",
            ],
            "annotations": [
                {"n": 1, "label": (250, 160), "target": (250, 160)},
                {"n": 2, "label": (720, 250), "target": (720, 250)},
                {"n": 3, "label": (240, 355), "target": (240, 355)},
                {"n": 4, "label": (780, 355), "target": (780, 355)},
                {"n": 5, "label": (190, 255), "target": (190, 255)},
            ],
            "speaker": [
                "ย้ำว่า machine กับ mold ใช้ model เดียวกันแต่บทบาทคนละแบบ",
                "ตรงนี้ต้องให้ทีมเข้าใจว่าต้นทุนเครื่องมาจาก field ฝั่ง machine ไม่ใช่ฝั่ง mold อย่างเดียว",
            ],
        },
        {
            "kind": "content",
            "title": "Workcenter Compatibility Matrix",
            "section": "Workcenter",
            "menu": "Manufacturing > Configuration > Work Centers > Compatibility Matrix",
            "image": "10b_workcenter_injection5_compatibility_matrix.png",
            "purpose": [
                "สอนฝั่ง machine ของ matrix ว่าเครื่องนี้รับ mold อะไรได้บ้าง",
                "เป็นแกนของการ auto suggest machine/mold ฝั่ง plastic",
            ],
            "clicks": [
                "1. เปิดแท็บ Compatibility Matrix ของเครื่อง",
                "2. ดูรายการ mold ที่เครื่องนี้ใช้ร่วมกันได้",
                "3. ดู usage และ limit เพื่อประเมินความเสี่ยงด้าน mold life",
                "4. ใช้ Add a line เมื่อต้องเพิ่ม compatibility ใหม่",
            ],
            "fields": [
                "Work Center ในตารางคือ mold ที่เครื่องนี้รองรับ",
                "Mold Status บอกสถานะ Normal, Warning หรือ Full",
                "Usage คือจำนวน shots ที่ใช้ไปแล้ว",
                "Limit คือ mold life limit ที่เอาไว้เตือนการซ่อมบำรุง",
            ],
            "annotations": [
                {"n": 1, "label": (180, 260), "target": (180, 260)},
                {"n": 2, "label": (190, 315), "target": (190, 315)},
                {"n": 3, "label": (905, 315), "target": (905, 315)},
                {"n": 4, "label": (95, 406), "target": (95, 406)},
            ],
            "speaker": [
                "สอนให้ทีมอ่านความสัมพันธ์ machine -> mold ให้เป็นจากฝั่งเครื่องก่อน",
                "ย้ำว่าถ้าไม่มี matrix ฝั่งนี้ ระบบ auto assign จะทำงานไม่ครบ",
            ],
        },
        {
            "kind": "content",
            "title": "Mold Master",
            "section": "Mold",
            "menu": "Manufacturing > Configuration > Work Centers (Mold)",
            "image": "11_mold_upper_w01_general.png",
            "purpose": [
                "สอน master data ของ mold ที่มีผลกับต้นทุนและอายุการใช้งาน",
                "ใช้หน้าจริงที่ตอนนี้โชว์ Is Mold? แล้วเพื่ออธิบาย field สำคัญทุกตัว",
            ],
            "clicks": [
                "1. เปิด mold record ที่ต้องการ",
                "2. ดู Is Mold? เพื่อยืนยันว่าระบบ treat record นี้เป็น mold จริง",
                "3. ดู Mold Cost / Hour, Cavities และ Mold Life Limit",
                "4. ดู Current Shots หรือ status เพื่อเช็กการใช้งานสะสม",
                "5. เปิด Mold Matrix เมื่อต้องดู product และ machine compatibility",
            ],
            "fields": [
                "Is Mold? ต้องถูกติ๊ก มิฉะนั้น logic mold จะไม่ทำงานครบ",
                "Mold Cost / Hour ใช้คำนวณต้นทุน mold",
                "Cavities ใช้ตีความ output ต่อหนึ่ง shot",
                "Mold Life Limit ใช้ควบคุมอายุ mold และเตือนบำรุงรักษา",
                "Current Shots และ Mold Status ใช้ติดตามการใช้งานสะสม",
            ],
            "annotations": [
                {"n": 1, "label": (250, 160), "target": (250, 160)},
                {"n": 2, "label": (725, 248), "target": (725, 248)},
                {"n": 3, "label": (768, 352), "target": (768, 352)},
                {"n": 4, "label": (760, 450), "target": (760, 450)},
                {"n": 5, "label": (170, 255), "target": (170, 255)},
            ],
            "speaker": [
                "อธิบายความต่างระหว่าง mold master กับ workcenter machine ให้ชัดอีกครั้ง",
                "ย้ำว่าหน้านี้มีผลทั้งด้าน planning, costing และ maintenance",
            ],
        },
        {
            "kind": "content",
            "title": "Mold Matrix",
            "section": "Mold",
            "menu": "Manufacturing > Configuration > Work Centers (Mold) > Mold Matrix",
            "image": "11b_mold_upper_w01_matrix.png",
            "purpose": [
                "สอนฝั่ง mold ของ matrix ซึ่งผูก Compatible Machines และ Produced Products Efficiency",
                "เป็นจุดที่ทำให้ระบบรู้ว่า mold ตัวนี้ใช้กับสินค้าอะไรและเครื่องไหนได้บ้าง",
            ],
            "clicks": [
                "1. เปิดแท็บ Mold Matrix บน mold",
                "2. ดู Compatible Machines ว่า mold ตัวนี้ลงเครื่องไหนได้",
                "3. ดู Produced Products Efficiency ว่าผลิตสินค้าอะไรได้บ้าง",
                "4. ดู Cycle Time และ Units / Hour เพื่อใช้วางแผนและอธิบาย std/actual performance",
            ],
            "fields": [
                "Compatible Machines คือรายชื่อเครื่องที่ใช้ mold ตัวนี้ได้",
                "Product คือสินค้าที่ mold นี้ผลิตได้",
                "Cycle Time (s) คือเวลามาตรฐานต่อรอบ",
                "Units / Hour ใช้แปลความสามารถในการผลิตต่อชั่วโมง",
                "สองตารางนี้รวมกันคือหัวใจของ auto mold assignment",
            ],
            "annotations": [
                {"n": 1, "label": (170, 280), "target": (170, 280)},
                {"n": 2, "label": (180, 340), "target": (180, 340)},
                {"n": 3, "label": (185, 550), "target": (185, 550)},
                {"n": 4, "label": (870, 550), "target": (870, 550)},
            ],
            "speaker": [
                "สอนให้ผู้เรียนเห็นว่าหน้านี้เป็นสะพานเชื่อม master data กับ execution จริง",
                "เชื่อมกับเรื่องที่เราแก้ logic ไปว่า mold เดียวไม่ควรถูก assign ซ้ำข้าม parallel workorders",
            ],
        },
        {
            "kind": "content",
            "title": "Transfer Plastic",
            "section": "Inventory",
            "menu": "Inventory > Operations > Transfers",
            "image": "12_transfer_plastic_gmp_trpl_00006.png",
            "purpose": [
                "สอน internal transfer ฝั่ง Plastic จากเอกสารจริง",
                "ใช้ยืนยันว่าระบบแยก operation type ฝั่ง plastic ออกจาก pharma แล้ว",
            ],
            "clicks": [
                "1. เปิด transfer ฝั่ง plastic",
                "2. ดู Operation Type ให้เป็น Transfer Plastic",
                "3. ดู Source Location และ Destination Location",
                "4. ดูรายการสินค้าใน operations lines",
                "5. ใช้ Validate เมื่อต้องยืนยัน movement",
            ],
            "fields": [
                "Operation Type ต้องเป็น Transfer Plastic",
                "Source/Destination Location ระบุทิศทางการย้าย stock",
                "Source Document ใช้ trace กลับไป MO ต้นทาง",
                "Demand และ Quantity ช่วยเทียบ requirement กับจำนวนที่ทำจริง",
            ],
            "annotations": [
                {"n": 1, "label": (120, 160), "target": (120, 160)},
                {"n": 2, "label": (235, 248), "target": (235, 248)},
                {"n": 3, "label": (875, 248), "target": (875, 248)},
                {"n": 4, "label": (245, 470), "target": (245, 470)},
                {"n": 5, "label": (77, 120), "target": (77, 120)},
            ],
            "speaker": [
                "ย้ำว่าต้องดู Operation Type ไม่ใช่ดูชื่อสินค้าอย่างเดียวเวลาบอกว่าฝั่งไหน",
                "สอนให้ inventory team ใช้ source document ย้อนกลับหา MO หรือ flow ต้นทางได้",
            ],
        },
        {
            "kind": "content",
            "title": "Transfer Pharma",
            "section": "Inventory",
            "menu": "Inventory > Operations > Transfers",
            "image": "13_transfer_pharma_gmp_trph_00001.png",
            "purpose": [
                "สอน internal transfer ฝั่ง Pharma และใช้เทียบกับ Plastic แบบเห็นภาพจริง",
                "ช่วยอธิบายให้ warehouse แยกเอกสารสองฝั่งได้ชัดเจน",
            ],
            "clicks": [
                "1. เปิด transfer ฝั่ง pharma",
                "2. ดู Operation Type ให้เป็น Transfer Pharma",
                "3. ดู Source Document ว่ามาจากงานหรือ MO ใด",
                "4. ดู operations lines เพื่อเช็กสินค้าที่กำลังย้าย",
            ],
            "fields": [
                "Operation Type ต้องเป็น Transfer Pharma",
                "Source Document อาจรวมหลาย MO หรือหลายงานที่ feed เข้ามาได้",
                "Operations table ใช้อ่านรายการสินค้าที่กำลังเคลื่อนย้าย",
                "Status ช่วยบอกว่างานพร้อม, done หรือยังต้องดำเนินการต่อ",
            ],
            "annotations": [
                {"n": 1, "label": (120, 160), "target": (120, 160)},
                {"n": 2, "label": (235, 248), "target": (235, 248)},
                {"n": 3, "label": (855, 248), "target": (855, 248)},
                {"n": 4, "label": (245, 470), "target": (245, 470)},
            ],
            "speaker": [
                "ใช้สไลด์นี้จับคู่กับ Transfer Plastic เพื่อให้ผู้เรียนเปรียบเทียบสองฝั่ง",
                "ย้ำว่าการ trace stock ข้ามโรงต้องดู transfer documents ให้เป็น",
            ],
        },
        {
            "kind": "content",
            "title": "Invoice & Accounting Endpoint",
            "section": "Accounting",
            "menu": "Accounting > Customers > Invoices",
            "image": "14_invoice_inv_d_26_04_00001.png",
            "purpose": [
                "สอนปลายทางด้านบัญชีของ flow หลัง sale และ delivery",
                "ทำให้ผู้เรียนเห็นว่าผลลัพธ์สุดท้ายของงานผลิตไม่ได้จบที่ MO แต่ไปถึง invoice และรายได้",
            ],
            "clicks": [
                "1. เปิด invoice จริงที่ออกจาก flow",
                "2. ดู Customer, Invoice Date และ Journal ก่อน",
                "3. ดู Invoice Lines เพื่ออธิบาย product, quantity, account และ taxes",
                "4. ดูยอดรวมด้านล่างเพื่อยืนยัน untaxed, VAT และ total",
                "5. ดูสถานะ Posted เพื่อยืนยันว่าพร้อมลงบัญชีจริงแล้ว",
            ],
            "fields": [
                "Customer, Invoice Date, Journal และ Payment Terms คือ header หลักของ invoice",
                "Invoice Lines ใช้ตรวจสินค้า ปริมาณ ราคา ภาษี และ account ที่ระบบผูกมา",
                "Analytic ใช้สำหรับ reporting หรือกระจายรายได้/ต้นทุน",
                "Untaxed, VAT และ Total ใช้ยืนยันยอดปลายทางของเอกสาร",
                "Posted คือสถานะพร้อมลงบัญชีจริง",
            ],
            "annotations": [
                {"n": 1, "label": (240, 165), "target": (240, 165)},
                {"n": 2, "label": (855, 245), "target": (855, 245)},
                {"n": 3, "label": (255, 470), "target": (255, 470)},
                {"n": 4, "label": (1230, 710), "target": (1230, 710)},
                {"n": 5, "label": (1225, 128), "target": (1225, 128)},
            ],
            "speaker": [
                "สอนให้ผู้เรียนเห็นว่าฝั่ง Accounting เป็นปลายทางของ flow ที่ผู้บริหารและ finance ใช้ตรวจผล",
                "ย้ำว่าถ้าต้อง trace ให้ครบ ต้องย้อนจาก invoice ไป sale, delivery, product และสุดท้ายถึง MO ได้",
            ],
        },
        {
            "kind": "closing",
            "title": "Closing Checklist",
            "items": [
                "เช็ก Product Routes และ Manufacturing Type ก่อนทุกครั้งที่ flow ไม่ไปต่อ",
                "เช็ก BoM และ child chain เมื่อ MO ไม่แตก component หรือ child MO ตามที่คาด",
                "เช็ก Reordering Rule, Buy route และ vendor เมื่อ shortage แล้ว PO ไม่เกิด",
                "เช็ก GMP Shop Floor, Work Order และ Scrap page เมื่อปัญหาอยู่ที่ execution หน้างาน",
                "เช็ก Workcenter / Mold Matrix เมื่อระบบเลือกเครื่องหรือ mold ไม่ถูก",
                "เช็ก Transfer Operation Type เพื่อยืนยันว่าของวิ่งฝั่ง Plastic หรือ Pharma ถูกต้อง",
                "เช็ก Invoice, Valuation และ Journal เมื่อจะตรวจผลทางบัญชีปลายทาง",
            ],
            "speaker": [
                "ปิดคลาสด้วย checklist ที่ใช้ได้จริงเวลาทีมเจอปัญหา",
                "ย้ำว่าก่อนแจ้ง bug ควรเช็ก route, bom, orderpoint, vendor และ mold mapping ก่อนเสมอ",
            ],
        },
    ]
)


def build_content(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, data["title"], data["section"])
    add_box(slide, Inches(0.25), Inches(0.95), Inches(4.9), Inches(1.35), "หน้าที่ของหน้า", data["purpose"], 10)
    add_box(slide, Inches(0.25), Inches(2.45), Inches(4.9), Inches(1.95), "ลำดับที่ต้องกด", data["clicks"], 9.5)
    add_box(slide, Inches(0.25), Inches(4.55), Inches(4.9), Inches(2.2), "Field และการทำงาน", data["fields"], 9.2)
    add_screenshot(slide, data["image"])
    add_footer(slide, data["section"], data["menu"])


def build_closing(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, data["title"], "What users should remember")
    add_box(slide, Inches(0.8), Inches(1.1), Inches(11.8), Inches(5.7), "Checklist", data["items"], 14)


def write_speaker_script():
    lines = []
    for idx, slide in enumerate(SLIDES, start=1):
        lines.append(f"Slide {idx}: {clean_text(slide['title'])}")
        lines.append("=" * (len(lines[-1])))
        if slide.get("menu"):
            lines.append(f"Menu path: {clean_text(slide['menu'])}")
        if slide.get("purpose"):
            lines.append("Purpose:")
            for item in slide["purpose"]:
                lines.append(f"- {clean_text(item)}")
        if slide.get("clicks"):
            lines.append("Click sequence:")
            for item in slide["clicks"]:
                lines.append(f"- {clean_text(item)}")
        if slide.get("fields"):
            lines.append("Field explanation:")
            for item in slide["fields"]:
                lines.append(f"- {clean_text(item)}")
        if slide.get("items"):
            lines.append("Items:")
            for item in slide["items"]:
                lines.append(f"- {clean_text(item)}")
        if slide.get("speaker"):
            lines.append("Speaker memory notes:")
            for item in slide["speaker"]:
                lines.append(f"- {clean_text(item)}")
        lines.append("")
    OUT_TXT.write_text("\n".join(lines), encoding="utf-8")


def main():
    ANNOTATED_DIR.mkdir(parents=True, exist_ok=True)
    for slide in SLIDES:
        if slide.get("kind") == "content":
            annotate_image(slide["image"], slide["annotations"])

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide in SLIDES:
        kind = slide["kind"]
        if kind == "title":
            build_title(prs, slide)
        elif kind == "agenda":
            build_agenda(prs, slide)
        elif kind == "closing":
            build_closing(prs, slide)
        else:
            build_content(prs, slide)

    prs.save(OUT_PPTX)
    write_speaker_script()


SLIDES.extend(
    [
        {
            "kind": "content",
            "title": "Purchase Order from Shortage",
            "section": "Purchase",
            "menu": "Purchase > Orders > Requests for Quotation",
            "image": "09_po_p00014.png",
            "purpose": [
                "ใช้สอนว่าของขาดไม่ได้จบที่ MO ทุกกรณี ถ้า component เป็น Buy route ต้องไป RFQ/PO",
                "หน้า PO ทำให้ทีม Purchase เข้าใจว่าตัวเองรับช่วงจาก planning อย่างไร",
            ],
            "clicks": [
                "1. เปิด RFQ หรือ PO ที่เกิดจาก shortage หรือ planning",
                "2. ดู Vendor และ Order Deadline ก่อนทุกครั้ง",
                "3. ดู Order Lines ว่าสินค้าตัวไหน quantity เท่าไรและ UoM อะไร",
                "4. ใช้ Confirm Order เมื่อต้องแปลงจาก RFQ เป็น PO จริง",
            ],
            "fields": [
                "Vendor คือผู้ขายที่ต้องมีใน master ไม่เช่นนั้น flow Buy จะไม่จบ",
                "Order Deadline และ Expected Arrival ใช้ประกอบการวางแผนรับของ",
                "Source Document ช่วย trace กลับไป shortage, requisition หรือ procurement",
                "Order Lines แสดงสินค้า ปริมาณ ราคา ภาษี และ analytic information",
            ],
            "annotations": [
                {"n": 1, "label": (250, 210), "target": (250, 210)},
                {"n": 2, "label": (760, 210), "target": (760, 210)},
                {"n": 3, "label": (245, 470), "target": (245, 470)},
                {"n": 4, "label": (65, 138), "target": (65, 138)},
            ],
            "speaker": [
                "ตรงนี้ตอบคำถามยอดฮิตว่า route ถูกแล้วแต่ทำไมไม่เกิด PO ซึ่งส่วนใหญ่เพราะไม่มี vendor หรือ source ไม่พร้อม",
                "ให้ผู้เรียนแยกให้ออกว่าหน้า PO เป็นปลายทางของ logic shortage ฝั่ง Buy route",
            ],
        },
        {
            "kind": "content",
            "title": "Promotion / FOC Sales Order",
            "section": "Flow 1",
            "menu": "Sales > Orders",
            "image": "06_sale_foc_sob_263069.png",
            "purpose": [
                "สอน flow ส่งเสริมการขายจาก SO จริงที่มี free line อยู่ใน UAT",
                "ใช้เป็นตัวอย่างว่า FOC ยังมี stock movement และ downstream documents",
            ],
            "clicks": [
                "1. เปิด SO ที่เป็น FOC flow",
                "2. ดู Delivery smart button เพื่อ trace ไป transfer และ delivery",
                "3. ดู Order Lines ให้เห็น line ปกติกับ line ฟรี [FREE]",
                "4. ดูยอดรวมด้านล่างเพื่ออธิบายผลทางการขายแม้มี free item",
            ],
            "fields": [
                "Customer, Invoice Address และ Delivery Address เป็นปลายทางของเอกสาร",
                "SO Type ใช้แยกประเภทการขายและ policy เฉพาะทาง",
                "Order Lines ต้องสังเกต line ฟรีที่ราคา 0 และ tag [FREE]",
                "Amount, VAT และ Total ใช้ตรวจว่ามูลค่าทางขายถูกคำนวณอย่างไร",
                "Delivery smart button เป็นจุดเชื่อมไป fulfillment",
            ],
            "annotations": [
                {"n": 1, "label": (120, 160), "target": (120, 160)},
                {"n": 2, "label": (760, 70), "target": (760, 70)},
                {"n": 3, "label": (222, 532), "target": (222, 532)},
                {"n": 4, "label": (1290, 720), "target": (1290, 720)},
            ],
            "speaker": [
                "ให้ผู้เรียนดูว่า free item อยู่ใน order เดียวกัน ไม่ได้แยกเอกสารพิเศษ",
                "ย้ำว่า FOC ไม่ได้แปลว่าไม่มี stock movement และยังต้อง trace delivery, invoice และ cost ต่อได้",
            ],
        },
        {
            "kind": "content",
            "title": "SO -> MTO Example",
            "section": "Flow 2",
            "menu": "Sales > Orders",
            "image": "07_sale_mto_sob_263070.png",
            "purpose": [
                "สอน demand ที่เริ่มจาก sale แล้ววิ่งต่อไป production",
                "ใช้เทียบกับ FOC เพื่อให้เห็นว่า UI คล้ายกันแต่ intent ของ flow ต่างกัน",
            ],
            "clicks": [
                "1. เปิด SO ที่ใช้เป็นตัวอย่าง MTO/MO downstream",
                "2. ดู Delivery smart button เพื่อเชื่อมไป picking ที่ระบบสร้างให้",
                "3. ดู Create Invoice เมื่ออธิบายขั้นตอนสร้าง invoice จาก sale",
                "4. ดู Order Lines และ delivered/invoiced quantities เพื่อเทียบความคืบหน้า",
            ],
            "fields": [
                "Delivery smart button ใช้ trace picking และ downstream inventory documents",
                "Create Invoice ใช้สร้าง invoice ตาม policy ของ flow นั้น",
                "Order Lines แสดง quantity ordered, delivered และ invoiced",
                "Statusbar ช่วยอธิบาย lifecycle จาก quotation ไป sales order",
            ],
            "annotations": [
                {"n": 1, "label": (120, 160), "target": (120, 160)},
                {"n": 2, "label": (760, 70), "target": (760, 70)},
                {"n": 3, "label": (70, 120), "target": (70, 120)},
                {"n": 4, "label": (255, 520), "target": (255, 520)},
            ],
            "speaker": [
                "อธิบายว่าถ้า setting ถูก ผู้ใช้ไม่ต้องเปิด MO ลูกเองจากหน้าขาย",
                "ให้ผู้เรียนจับความสัมพันธ์ระหว่าง sale line, delivery และ production trace ต่อไป",
            ],
        },
        {
            "kind": "content",
            "title": "Manufacturing Order from Replenishment",
            "section": "Flow 3",
            "menu": "Manufacturing > Operations > Manufacturing Orders",
            "image": "08_mo_mts_gmp_moph_00011.png",
            "purpose": [
                "สอน MO ที่เกิดจาก replenishment/Min-Max ไม่ได้เกิดจาก SO โดยตรง",
                "ใช้เป็นจุดรวมในการอธิบาย component readiness, workorders, transfers และ cost",
            ],
            "clicks": [
                "1. เปิด MO ที่เกิดจาก replenishment",
                "2. ดู Product, Quantity To Produce และ BoM ที่ใช้",
                "3. ดู Component Status ว่าวัตถุดิบพร้อมหรือไม่",
                "4. ใช้ smart buttons เช่น Transfers, Shop Floor, MO Cost และ Workorders เพื่อ drill down",
                "5. เปิด Components tab เมื่อต้องสอน source location และ quantity to consume",
            ],
            "fields": [
                "Product และ Quantity To Produce คือสิ่งที่จะผลิตจริง",
                "BoM ชี้กลับไป master data ต้นทางของงานนี้",
                "Scheduled Date และ Scheduled End ใช้ประกอบ planning",
                "Component Status ช่วยประเมิน readiness ก่อนเริ่มงานจริง",
                "Components tab ใช้ตรวจ requirement รายชิ้นและ source location",
            ],
            "annotations": [
                {"n": 1, "label": (145, 160), "target": (145, 160)},
                {"n": 2, "label": (245, 250), "target": (245, 250)},
                {"n": 3, "label": (815, 250), "target": (815, 250)},
                {"n": 4, "label": (540, 70), "target": (540, 70)},
                {"n": 5, "label": (165, 425), "target": (165, 425)},
            ],
            "speaker": [
                "เน้นคำว่า replenishment ให้ชัด เพราะนี่คือหัวใจของ MTS/Min-Max flow",
                "สอนให้ผู้ใช้เช็ก Component Status ก่อนเสมอ ไม่เช่นนั้นเริ่มงานไปแล้วอาจติดวัตถุดิบไม่พร้อม",
            ],
        },
        {
            "kind": "content",
            "title": "GMP Shop Floor Dashboard",
            "section": "Execution",
            "menu": "GMP Shop Floor > Dashboard",
            "image": "16_gmp_shop_floor_dashboard.png",
            "purpose": [
                "สอนหน้ารวมงานหน้างานจริงที่หัวหน้างานใช้เปิดทุกวัน",
                "ทำให้ผู้เรียนเห็น card-based execution และจุดเข้า console ของแต่ละ MO",
            ],
            "clicks": [
                "1. เข้า GMP Shop Floor จากหน้า Home",
                "2. ใช้แท็บ Dashboard, Plastic Shop Floor และ Pharma Shop Floor เพื่อแยกมุมมองงาน",
                "3. ใช้ Search และ filter ด้านบนเพื่อหางานที่ต้องการ",
                "4. กด Open Console บน card ที่ต้องการเริ่มทำงาน",
            ],
            "fields": [
                "แต่ละ card แสดง MO, origin และชื่อสินค้าที่กำลังผลิต",
                "เป้าหมายคือ quantity target ของงานนั้น",
                "Tasks แยกเป็น Ready, In Progress และ Done เพื่อดูความคืบหน้า",
                "Open Console คือจุดเข้า execution ของหน้างานจริง",
            ],
            "annotations": [
                {"n": 1, "label": (88, 36), "target": (88, 36)},
                {"n": 2, "label": (251, 40), "target": (251, 40)},
                {"n": 3, "label": (240, 95), "target": (240, 95)},
                {"n": 4, "label": (241, 461), "target": (241, 461)},
            ],
            "speaker": [
                "นี่คือหน้า operational dashboard ของจริง ไม่ใช่ master data",
                "สอนให้หัวหน้างานอ่าน card จากเลข MO, product, เป้าหมาย และ tasks ก่อนกด Open Console",
            ],
        },
        {
            "kind": "content",
            "title": "Work Order / Shop Floor",
            "section": "Execution",
            "menu": "Manufacturing > Operations > Work Orders",
            "image": "15_workorder_operation_112.png",
            "purpose": [
                "สอนหน้ารายละเอียด work order ที่ operator หรือ supervisor ใช้ดูการทำงานจริง",
                "เป็นจุดที่เชื่อม machine, mold, quantity, time tracking และเอกสารแม่เข้าด้วยกัน",
            ],
            "clicks": [
                "1. เปิด work order ที่ต้องการ",
                "2. ดู Work Center และ Product ว่า operation นี้ทำที่เครื่องไหนและผลิตอะไร",
                "3. ดู Molds และ Mold Cost เมื่อเป็นงานฝั่ง plastic",
                "4. ใช้ Manufacturing Order link ย้อนกลับไปหาเอกสารแม่",
                "5. ดู Time Tracking เมื่อต้องอธิบายเวลาเริ่มงานและ productivity",
            ],
            "fields": [
                "Work Center บอกว่า operation นี้ทำที่เครื่องใด",
                "Product และ Quantity บอกว่ากำลังผลิตอะไรจำนวนเท่าไร",
                "Molds และ Mold Cost สำคัญกับงาน plastic และ costing",
                "Expected Duration คือเวลามาตรฐานที่ระบบคำนวณให้",
                "Time Tracking ใช้ดู employee productivity, start/end และเวลาจริง",
            ],
            "annotations": [
                {"n": 1, "label": (210, 240), "target": (210, 240)},
                {"n": 2, "label": (230, 160), "target": (230, 160)},
                {"n": 3, "label": (1025, 245), "target": (1025, 245)},
                {"n": 4, "label": (1140, 160), "target": (1140, 160)},
                {"n": 5, "label": (865, 495), "target": (865, 495)},
            ],
            "speaker": [
                "ย้ำว่าหน้านี้คือ operational truth ของหน้างาน ไม่ใช่แค่ข้อมูลสรุป",
                "เชื่อมให้ผู้เรียนเห็นว่าข้อมูล mold และเวลาในหน้านี้มีผลต่อ actual cost ด้วย",
            ],
        },
        {
            "kind": "content",
            "title": "Scrap Transaction",
            "section": "Execution",
            "menu": "Inventory > Operations > Scrap",
            "image": "17_scrap_training.png",
            "purpose": [
                "ใช้เอกสาร scrap จริงที่สร้างบน UAT เพื่อให้ training มี execution page ครบ",
                "สอนว่าของเสียกระทบ stock และเป็นเอกสารปฏิบัติการที่ต้อง trace ได้",
            ],
            "clicks": [
                "1. เปิดหน้าสร้าง Scrap หรือเอกสาร scrap ที่มีอยู่",
                "2. เลือก Product และ Quantity ที่ต้อง scrap",
                "3. เลือก Source Location และ Scrap Location ให้ถูก",
                "4. ระบุ Source Document หรือ Scrap Reason ถ้ามี",
                "5. กด Validate เพื่อยืนยัน movement ของ scrap",
            ],
            "fields": [
                "Product คือสินค้าที่จะ scrap",
                "Quantity ต้องเป็นจำนวนที่ต้องตัดออกจาก stock จริง",
                "Source Location คือที่ stock อยู่ก่อน scrap",
                "Scrap Location คือที่ปลายทางของของเสีย",
                "Source Document ใช้เชื่อมกับงานหรือเหตุผลที่เกี่ยวข้อง",
                "Replenish Quantities ใช้เมื่อองค์กรต้องการให้ scrap ไปกระตุ้น procurement เพิ่ม",
            ],
            "annotations": [
                {"n": 1, "label": (240, 245), "target": (240, 245)},
                {"n": 2, "label": (235, 278), "target": (235, 278)},
                {"n": 3, "label": (850, 245), "target": (850, 245)},
                {"n": 4, "label": (845, 308), "target": (845, 308)},
                {"n": 5, "label": (46, 137), "target": (46, 137)},
            ],
            "speaker": [
                "บอกผู้เรียนว่านี่คือ transaction จริงที่สร้างบน UAT เพื่อใช้ประกอบ training",
                "ย้ำว่า scrap ไม่ใช่แค่บันทึกของเสีย แต่เป็น stock movement ที่กระทบ inventory จริง",
            ],
        },
    ]
)


def _font(size):
    candidates = [
        Path(r"C:\Windows\Fonts\tahoma.ttf"),
        Path(r"C:\Windows\Fonts\arial.ttf"),
    ]
    for path in candidates:
        if path.exists():
            return ImageFont.truetype(str(path), size=size)
    return ImageFont.load_default()


def annotate_image(src_name, annotations):
    src = SHOT_DIR / src_name
    dst = ANNOTATED_DIR / src_name
    dst.parent.mkdir(parents=True, exist_ok=True)
    img = Image.open(src).convert("RGBA")
    draw = ImageDraw.Draw(img)
    num_font = _font(26)

    for item in annotations:
        label = item["label"]
        target = item["target"]
        number = str(item["n"])
        draw.line([label, target], fill=(198, 40, 40, 255), width=6)
        tx, ty = target
        lx, ly = label
        draw.polygon([(tx, ty), (tx - 14, ty - 7), (tx - 14, ty + 7)], fill=(198, 40, 40, 255))
        r = 26
        draw.ellipse((lx - r, ly - r, lx + r, ly + r), fill=(198, 40, 40, 255), outline=(255, 255, 255, 255), width=4)
        bbox = draw.textbbox((0, 0), number, font=num_font)
        tw = bbox[2] - bbox[0]
        th = bbox[3] - bbox[1]
        draw.text((lx - tw / 2, ly - th / 2 - 2), number, font=num_font, fill=(255, 255, 255, 255))

    img.save(dst)
    return dst


def set_run_style(run, size, bold=False, color=TEXT):
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def set_para(paragraph, text, size, bold=False, color=TEXT):
    paragraph.clear()
    run = paragraph.add_run()
    run.text = clean_text(text)
    set_run_style(run, size, bold, color)


def add_banner(slide, title, subtitle):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.62))
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()

    box = slide.shapes.add_textbox(Inches(0.3), Inches(0.1), Inches(9.0), Inches(0.3))
    set_para(box.text_frame.paragraphs[0], title, 24, True, WHITE)

    if subtitle:
        sub = slide.shapes.add_textbox(Inches(9.5), Inches(0.12), Inches(3.4), Inches(0.22))
        tf = sub.text_frame
        tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
        set_para(tf.paragraphs[0], subtitle, 10, False, WHITE)


def add_box(slide, left, top, width, height, title, bullets, font_size=10):
    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = WHITE
    card.line.color.rgb = LINE
    card.line.width = Pt(1)

    title_box = slide.shapes.add_textbox(left + Inches(0.12), top + Inches(0.08), width - Inches(0.24), Inches(0.22))
    set_para(title_box.text_frame.paragraphs[0], title, 12, True, GREEN_DARK)

    body_box = slide.shapes.add_textbox(left + Inches(0.14), top + Inches(0.34), width - Inches(0.28), height - Inches(0.42))
    tf = body_box.text_frame
    tf.word_wrap = True
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    for idx, bullet in enumerate(bullets):
        para = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        para.bullet = True
        para.space_after = Pt(2)
        para.line_spacing = 1.05
        run = para.add_run()
        run.text = fill(clean_text(bullet), 42)
        set_run_style(run, font_size, False, TEXT)


def add_screenshot(slide, image_name):
    path = ANNOTATED_DIR / image_name
    frame = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(5.4), Inches(0.95), Inches(7.55), Inches(5.85))
    frame.fill.solid()
    frame.fill.fore_color.rgb = WHITE
    frame.line.color.rgb = LINE
    frame.line.width = Pt(1)
    slide.shapes.add_picture(str(path), Inches(5.52), Inches(1.07), width=Inches(7.31), height=Inches(5.6))


def add_footer(slide, section, menu):
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0), Inches(7.08), Inches(13.333), Inches(0.42))
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN_DARK
    band.line.fill.background()

    left = slide.shapes.add_textbox(Inches(0.3), Inches(7.16), Inches(2.7), Inches(0.18))
    set_para(left.text_frame.paragraphs[0], section, 10, True, WHITE)
    right = slide.shapes.add_textbox(Inches(3.1), Inches(7.16), Inches(9.8), Inches(0.18))
    right_tf = right.text_frame
    right_tf.paragraphs[0].alignment = PP_ALIGN.RIGHT
    set_para(right_tf.paragraphs[0], menu, 10, False, WHITE)


def build_title(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    hero = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.75), Inches(1.0), Inches(11.8), Inches(5.3))
    hero.fill.solid()
    hero.fill.fore_color.rgb = WHITE
    hero.line.color.rgb = LINE
    band = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.RECTANGLE, Inches(0.75), Inches(1.0), Inches(11.8), Inches(0.8))
    band.fill.solid()
    band.fill.fore_color.rgb = GREEN
    band.line.fill.background()
    title = slide.shapes.add_textbox(Inches(1.0), Inches(2.0), Inches(10.5), Inches(1.0))
    set_para(title.text_frame.paragraphs[0], data["title"], 28, True, GREEN_DARK)
    sub = slide.shapes.add_textbox(Inches(1.0), Inches(2.95), Inches(7.2), Inches(1.5))
    tf = sub.text_frame
    first = True
    for line in data["subtitle"].splitlines():
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        set_para(p, line, 16, False, MUTED)
    box = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.0), Inches(5.05), Inches(4.0), Inches(0.55))
    box.fill.solid()
    box.fill.fore_color.rgb = RGBColor(224, 245, 229)
    box.line.fill.background()
    caption = slide.shapes.add_textbox(Inches(1.15), Inches(5.2), Inches(3.7), Inches(0.18))
    set_para(caption.text_frame.paragraphs[0], "Real screenshots and real UAT transactions", 11, True, GREEN_DARK)


def build_agenda(prs, data):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG
    add_banner(slide, data["title"], "Detailed session map")
    add_box(slide, Inches(0.8), Inches(1.1), Inches(11.8), Inches(5.7), "Agenda", data["items"], 16)


if __name__ == "__main__":
    main()
