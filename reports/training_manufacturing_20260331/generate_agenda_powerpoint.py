from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt


OUT = Path(
    r"c:\365_project\TheCool18e\Dev\reports\training_manufacturing_20260331\manufacturing_training_agenda_20260331.pptx"
)


AGENDA_SECTIONS = [
    {
        "title": "Opening และภาพรวมระบบ",
        "time": "00:00-00:15",
        "points": [
            "อธิบายเป้าหมายของระบบ Manufacturing และภาพรวมโรงงาน Plastic กับ Pharma",
            "ปู flow ใหญ่ของบริษัทจาก Demand ไปจนถึง Accounting",
            "เปิดหน้า Home และ Manufacturing เพื่อให้ทุกคนเห็นภาพรวมเมนูที่เกี่ยวข้อง",
            "ผลลัพธ์ที่ต้องได้คือผู้เข้าอบรมเห็นภาพ end-to-end ก่อนลงรายละเอียด",
        ],
    },
    {
        "title": "Settings และ Master Data",
        "time": "00:15-00:35",
        "points": [
            "สอน Product Routes, Manufacturing Type, BoM, BoM Picking Type, Orderpoint, Vendor และ Mold Mapping",
            "เปิด Product และ BoM ของ FG-PSS-TH-01005 เป็นตัวอย่างจริง",
            "ย้ำว่า setting กลุ่มนี้เป็นต้นเหตุหลักของ flow ที่ไม่ auto",
            "ผลลัพธ์ที่ต้องได้คือผู้เข้าอบรมรู้ว่าต้องเช็กอะไรเมื่อระบบไม่ไปต่อ",
        ],
    },
    {
        "title": "Flow 1: Promotion / SO 0 บาท",
        "time": "00:35-01:00",
        "points": [
            "อธิบาย free item และ FOC logic ตั้งแต่ quotation ไปจนถึง accounting",
            "เปิด S11563 และ SOB-263069 เพื่อ trace SO -> MO -> Delivery -> Invoice",
            "ชี้ให้เห็นว่า SO 0 บาท ยังมี stock movement และ cost",
            "ผลลัพธ์ที่ต้องได้คือผู้เข้าอบรม trace FOC chain ได้ด้วยตัวเอง",
        ],
    },
    {
        "title": "Flow 2: SO -> MTO",
        "time": "01:00-01:25",
        "points": [
            "อธิบาย logic ของ demand จาก sale และ procurement group",
            "เปิด SOB-263070 และ child MO ที่เกี่ยวข้อง",
            "ชี้ให้เห็นความต่างระหว่าง make_to_order กับ make_to_stock",
            "ผลลัพธ์ที่ต้องได้คือผู้เข้าอบรมแยกได้ว่า case ไหนควรสร้าง MO จาก SO",
        ],
    },
    {
        "title": "Flow 3: MTS / Min-Max",
        "time": "01:25-01:45",
        "points": [
            "อธิบาย orderpoint, reordering rule, กรณีของพอ และกรณี stock ต่ำกว่า min",
            "เปิด FG-PNC-TH-01001 และ GMP/MOPH/00011 เป็นตัวอย่างจริงของ replenishment",
            "ใช้ FG-PSS-TH-01005 เป็นตัวอย่างของ ORDERPOINT_ONLY",
            "ผลลัพธ์ที่ต้องได้คือผู้เข้าอบรมตอบได้ว่าของขาดแล้วควรไป MO หรือ PO เพราะอะไร",
        ],
    },
    {
        "title": "Shopfloor, Workorder, Scrap",
        "time": "01:45-02:10",
        "points": [
            "สอน start, good qty, reject qty, done และ scrap",
            "เปิด workorder ของ GMP/MOPL/00014 และ scrap SP/00011",
            "อธิบายผลของ qty และ scrap ต่อ stock และ cost",
            "ผลลัพธ์ที่ต้องได้คือ operator และ supervisor เข้าใจจุดที่ต้องกรอกและต้องตรวจ",
        ],
    },
    {
        "title": "Mold Management",
        "time": "02:10-02:30",
        "points": [
            "สอน mold map กับ product และ workcenter, auto assign mold และ mold life",
            "เปิด GMP/MOPL/00014 และ mold ที่ใช้กับ SM-PLS-UP-01001",
            "ย้ำ logic ใหม่ว่า mold เดียวไม่ถูก assign ซ้ำข้าม parallel workorder",
            "ผลลัพธ์ที่ต้องได้คือผู้เข้าอบรมรู้ว่าถ้า mold ไม่ขึ้นต้องเช็กอะไร",
        ],
    },
    {
        "title": "Transfer, Inventory, Costing",
        "time": "02:30-02:50",
        "points": [
            "อธิบาย Transfer Plastic กับ Transfer Pharma และ semi ข้ามโรง",
            "สรุป raw, machine, labor, mold cost และความต่างระหว่าง std กับ actual",
            "ชี้เอกสารที่ใช้ trace cost ที่ปลายทาง",
            "ผลลัพธ์ที่ต้องได้คือทีม Production, Warehouse และ Accounting คุยบนข้อมูลชุดเดียวกัน",
        ],
    },
    {
        "title": "Q&A และ Checklist",
        "time": "02:50-03:00",
        "points": [
            "recap จุดสำคัญของทั้ง session",
            "สรุป checklist ก่อนใช้งานจริงและคำถามที่พบบ่อย",
            "ชี้ next step หลัง training ว่าต้องไป UAT อะไรต่อ",
            "ผลลัพธ์ที่ต้องได้คือ key user พร้อมไป test ต่อได้จริง",
        ],
    },
]


def style_run(run, size, bold=False, color=None):
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color


def add_header_band(slide):
    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.333), Inches(0.45))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(22, 101, 52)
    band.line.fill.background()


def add_title_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)
    add_header_band(slide)

    slide.shapes.title.text = "Manufacturing Training Agenda"
    for p in slide.shapes.title.text_frame.paragraphs:
        for run in p.runs:
            style_run(run, 28, bold=True, color=RGBColor(31, 41, 55))

    subtitle = slide.placeholders[1]
    subtitle.text = "Odoo 18 Enterprise\nSession 3 ชั่วโมง\nPlastic + Pharma Integrated Flow"
    for p in subtitle.text_frame.paragraphs:
        for run in p.runs:
            style_run(run, 16, color=RGBColor(75, 85, 99))


def add_objective_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_header_band(slide)

    slide.shapes.title.text = "Training Objective"
    for p in slide.shapes.title.text_frame.paragraphs:
        for run in p.runs:
            style_run(run, 24, bold=True, color=RGBColor(22, 101, 52))

    body = slide.placeholders[1].text_frame
    body.clear()
    bullets = [
        "เข้าใจภาพรวมระบบ Manufacturing ของบริษัทและความต่างของ Plastic กับ Pharma",
        "เข้าใจ Promotion, MTO, MTS / Min-Max และจุดที่ระบบ auto ให้",
        "เข้าใจ Shopfloor, Scrap, Mold, Transfer และ Costing ในระดับใช้งานจริง",
        "รู้ว่าถ้าระบบไม่ไปต่อ ต้องเช็ก route, BoM, orderpoint, vendor และ mold mapping ก่อน",
    ]
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        for run in p.runs:
            style_run(run, 18, color=RGBColor(31, 41, 55))


def add_agenda_slide(prs, title, time_range, points):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_header_band(slide)

    slide.shapes.title.text = f"{time_range}  {title}"
    for p in slide.shapes.title.text_frame.paragraphs:
        for run in p.runs:
            style_run(run, 22, bold=True, color=RGBColor(22, 101, 52))

    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, point in enumerate(points):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = point
        p.level = 0
        for run in p.runs:
            style_run(run, 16, color=RGBColor(31, 41, 55))


def add_summary_slide(prs):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)
    add_header_band(slide)

    slide.shapes.title.text = "Expected Outcome"
    for p in slide.shapes.title.text_frame.paragraphs:
        for run in p.runs:
            style_run(run, 24, bold=True, color=RGBColor(22, 101, 52))

    body = slide.placeholders[1].text_frame
    body.clear()
    bullets = [
        "ทีมแยก Promotion, MTO และ MTS ได้",
        "ทีม trace SO -> MO -> Transfer -> Delivery -> Accounting ได้",
        "ทีมใช้งาน Shopfloor และ Mold ได้ในระดับ operation จริง",
        "ทีมรู้ว่าถ้าระบบไม่ไปต่อ ต้องเช็กอะไรและต้อง escalate ให้ใคร",
    ]
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        for run in p.runs:
            style_run(run, 18, color=RGBColor(31, 41, 55))


def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    add_title_slide(prs)
    add_objective_slide(prs)
    for section in AGENDA_SECTIONS:
        add_agenda_slide(prs, section["title"], section["time"], section["points"])
    add_summary_slide(prs)

    prs.save(OUT)
    print(f"saved {OUT}")


if __name__ == "__main__":
    main()
