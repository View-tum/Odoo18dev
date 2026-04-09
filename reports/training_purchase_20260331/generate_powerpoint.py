from pathlib import Path
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Output path
OUT = Path(r"c:\365_project\TheCool18e\Dev\reports\training_purchase_20260331\Odoo18_Purchase_3hr_Intensive.pptx")

SLIDES = [
    {
        "layout": "title",
        "title": "Odoo 18 Purchase: 3-Hour Intensive",
        "subtitle": "Odoo 18 Enterprise - TheCool 18e\nWorkshop เจาะลึก 4 Flow สำคัญในเวลาจำกัด\nหลักสูตรเร่งด่วนสำหรับผู้ใช้งาน",
    },
    {
        "title": "ตารางการอบรม (Time-Blocked)",
        "bullets": [
            "09:00 - 09:30: พื้นฐานจัดซื้อ และการตั้งค่า Master Data ที่จำเป็น",
            "09:30 - 10:00: Flow 1 - การซื้อวัตถุดิบ (RM) และการจัดการของเสีย (Scrap)",
            "10:00 - 10:30: Flow 2 - การซื้อทรัพย์สิน (Assets) และการผูกความสัมพันธ์",
            "10:30 - 11:00: Flow 3 - สินค้าสิ้นเปลือง (Consumables) และการคุมแผนก",
            "11:00 - 11:30: Flow 4 - งานบริการ (Services) และการรับงาน (Service Entry)",
            "11:30 - 12:00: สรุปบัญชี E2E สรุปประเด็น และถาม-ตอบ",
        ],
    },
    {
        "title": "Part 1: Overview & Setting (09:00 - 09:30)",
        "bullets": [
            "1. Master Data: Vendor (Payment Terms) และ Product Category (Account)",
            "2. Internal Control: ทุกการสั่งซื้อต้องเริ่มที่ใบขอซื้อ (PR) เพื่อคุมงบ",
            "3. Workflow Approval: อนุมัติผ่านระบบ oi_workflow (Dept -> Finance -> MD)",
            "4. Reordering Rule: ตั้งจุดสั่งซื้อเพื่อลดงานคีย์ PR ด้วยมือ",
        ],
        "footer": "Setting Need: Foundation for accurate Data Flow.",
    },
    {
        "title": "Flow 1: 📦 วัตถุดิบ (RM) + Scrap (09:30 - 10:00)",
        "bullets": [
            "1. PR Creation: ระบุสินค้าเม็ดพลาสติก 100 kg พร้อมงบประมาณ",
            "2. PO & Receipt: ยืนยันราคาและรับของเข้าคลังพร้อมระบุ Lot Number",
            "3. Scrap Order: บันทึกของเสีย 5 kg เพื่อยอดสต็อกที่แม่นยำ",
            "4. Landed Cost: ปันส่วนมูลค่า 5kg กลับไปเป็นต้นทุนให้ 95kg ที่เหลือทันที!",
        ],
        "footer": "Accounting Tip: Landed Cost will affect Total Stock Valuation.",
    },
    {
        "title": "Flow 2: 🛠️ ทรัพย์สิน (Assets) แม่-ลูก (10:00 - 10:30)",
        "bullets": [
            "1. CAPEX Budget: เปิด PR ซื้อเครื่องจักรหลักและส่วนประกอบลูกในใบเดียว",
            "2. Asset Receipt: รับเข้าสถานที่ Asset Location เพื่อคุมทะเบียนทรัพย์สิน",
            "3. Vendor Bill: บันทึกบัญชีเข้า 'สินทรัพย์ระหว่างติดตั้ง' เพื่อรอ Post",
            "4. Asset Hierarchy: ผูก Parent Asset เพื่อบริหารมูลค่าทรัพย์สินรวบทั้งกลุ่ม",
        ],
        "footer": "Asset Automation: Automatic ID Creation on Bill Validation.",
    },
    {
        "title": "Flow 3: 🧤 สินค้าสิ้นเปลือง (Consumables) (10:30 - 11:00)",
        "bullets": [
            "เป้าหมาย: รับของนับจำนวนได้ (On Hand) แต่ตัดจ่ายทันที (Expense)",
            "1. Analytic Rule: ต้องระบุ Analytic Account ของแผนกที่เบิกใช้ใน PR",
            "2. Manual Valuation: รับเข้าสต็อกเพื่อนับจำนวน แต่บัญชีสะท้อน $0 ทันที",
            "3. Internal Transfer: ใช้ใบโอนย้ายภายใน เพื่อตัดจำนวนออกจากคลังจริง",
        ],
        "footer": "Controlled Items: PPE, Gloves, Office Stationery.",
    },
    {
        "title": "Flow 4: 📞 งานบริการ (Services) (11:00 - 11:30)",
        "bullets": [
            "เป้าหมาย: จ้างเหมาบริการและบันทึกรับงาน (Service Entry)",
            "1. PO Action: บันทึกรายการบริการซ่อมบำรุง/จ้างเหมาในหน้าใบสั่งซื้อ",
            "2. Service Entry: แก้ไขปริมาณ Received จาก 0 เป็น 1 (ได้รับการบริการแล้ว)",
            "3. Accounting Bill: สร้างใบแจ้งหนี้ตั้งเจ้าหนี้ตามสัดส่วนงานที่ได้รับจริง",
        ],
        "footer": "Note: Services do not generate Warehouse Receipts.",
    },
    {
        "title": "Part 3: บัญชี E2E และสรุป (11:30 - 12:00)",
        "bullets": [
            "Receipt (GRN): [Dr. สต็อก / Cr. สินค้าระหว่างทาง (Accrued)]",
            "Vendor Bill: [Dr. สินค้าระหว่างทาง / Cr. เจ้าหนี้การค้า]",
            "Audit Trail: ใช้ Smart Button ในหน้า PO เพื่อดูความเชื่อมโยงเอกสาร",
            "Q&A: ถาม-ตอบปัญหาหน้างาน และสรุปแนวทางการใช้ UAT ใน DB11",
        ],
    },
    {
        "title": "สรุป (Final Wrap-up)",
        "bullets": [
            "มาตรฐานจัดซื้อ คือจุดเริ่มต้นของกำไรบริษัท",
            "วินัยในการทำ Step-by-step สำคัญกว่าความเร็ว",
            "ขอบคุณผู้เข้าร่วมการอบรมทุกท่านครับ",
        ],
    },
]

def set_run_font(run, size, bold=False, color=None):
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    if color:
        run.font.color.rgb = color

def add_footer(slide, text):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(6.75), Inches(12.3), Inches(0.35))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = text
    set_run_font(run, 10, color=RGBColor(102, 102, 102))

def add_title_slide(prs, title, subtitle):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(245, 247, 250)

    title_shape = slide.shapes.title
    title_shape.text = title
    for p in title_shape.text_frame.paragraphs:
        for run in p.runs:
            set_run_font(run, 28, bold=True, color=RGBColor(31, 41, 55))

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    for p in subtitle_shape.text_frame.paragraphs:
        for run in p.runs:
            set_run_font(run, 16, color=RGBColor(75, 85, 99))

    band = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(0.45))
    band.fill.solid()
    band.fill.fore_color.rgb = RGBColor(22, 101, 52)
    band.line.fill.background()

def add_content_slide(prs, title, bullets, footer=None):
    slide = prs.slides.add_slide(prs.slide_layouts[1])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(255, 255, 255)

    title_shape = slide.shapes.title
    title_shape.text = title
    for p in title_shape.text_frame.paragraphs:
        for run in p.runs:
            set_run_font(run, 24, bold=True, color=RGBColor(22, 101, 52))

    body = slide.placeholders[1].text_frame
    body.clear()
    for idx, bullet in enumerate(bullets):
        p = body.paragraphs[0] if idx == 0 else body.add_paragraph()
        p.text = bullet
        p.level = 0
        for run in p.runs:
            set_run_font(run, 18, color=RGBColor(31, 41, 55))

    accent = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0.4), Inches(0.8), Inches(0.18), Inches(5.7))
    accent.fill.solid()
    accent.fill.fore_color.rgb = RGBColor(22, 101, 52)
    accent.line.fill.background()

    if footer:
        add_footer(slide, footer)

def main():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for slide_data in SLIDES:
        if slide_data.get("layout") == "title":
            add_title_slide(prs, slide_data["title"], slide_data["subtitle"])
        else:
            add_content_slide(
                prs,
                slide_data["title"],
                slide_data["bullets"],
                slide_data.get("footer"),
            )

    prs.save(OUT)
    print(f"saved {OUT}")

if __name__ == "__main__":
    main()
