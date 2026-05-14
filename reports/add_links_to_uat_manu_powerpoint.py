from __future__ import annotations

import re
from collections import OrderedDict
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_AUTO_SIZE
from pptx.dml.color import RGBColor
from pptx.util import Inches, Pt


SOURCE_XLSX = Path(r"C:\Users\tumsu\Downloads\UAT_MANU.xlsx")
INPUT_PPTX = Path(r"C:\Users\tumsu\Downloads\GMP_Manufacturing_UAT_from_UAT_MANU.pptx")
OUTPUT_PPTX = Path(r"C:\Users\tumsu\Downloads\GMP_Manufacturing_UAT_from_UAT_MANU.pptx")
BACKUP_PPTX = Path(r"C:\Users\tumsu\Downloads\GMP_Manufacturing_UAT_from_UAT_MANU_before_links.pptx")

ACTUAL_SHEET_PREFIX = tuple(f"{i:02d}_" for i in range(1, 10))

NAV_BLUE = RGBColor(31, 78, 121)
NAV_GOLD = RGBColor(255, 210, 66)
TEXT_WHITE = RGBColor(255, 255, 255)
TEXT_DARK = RGBColor(18, 34, 56)


def norm(value) -> str:
    return re.sub(r"\s+", " ", str(value or "").strip())


def get_headers(ws) -> dict[str, int]:
    headers: dict[str, int] = {}
    for col in range(1, ws.max_column + 1):
        value = norm(ws.cell(6, col).value)
        if value:
            headers[value] = col
    return headers


def collect_cases() -> list[dict]:
    wb = load_workbook(SOURCE_XLSX, data_only=False)
    cases: list[dict] = []
    for ws in wb.worksheets:
        if not ws.title.startswith(ACTUAL_SHEET_PREFIX):
            continue
        headers = get_headers(ws)
        if "Case ID" not in headers:
            continue
        scenario_col = headers.get("Scenario ทดสอบ") or headers.get("Scenario à¸—à¸”à¸ªà¸­à¸š") or 5
        for row in range(7, ws.max_row + 1):
            case_id = norm(ws.cell(row, headers["Case ID"]).value)
            if not re.match(r"^MU\d{2}-\d{2}$", case_id):
                continue
            cases.append(
                {
                    "case_id": case_id,
                    "sheet": ws.title,
                    "row": row,
                    "section": norm(ws["B2"].value) or ws.title,
                    "scenario": norm(ws.cell(row, scenario_col).value),
                }
            )
    return cases


def move_slide(prs: Presentation, old_index: int, new_index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[old_index]
    slide_id_list.remove(slide_id)
    slide_id_list.insert(new_index, slide_id)


def write_text(shape, text: str, size: float, bold: bool = False, color: RGBColor = TEXT_DARK) -> None:
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    run = p.add_run()
    run.text = text
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.color.rgb = color


def add_button(
    slide,
    text: str,
    left,
    top,
    width,
    height,
    *,
    target_slide=None,
    hyperlink: str | None = None,
    fill: RGBColor = NAV_BLUE,
    text_color: RGBColor = TEXT_WHITE,
):
    shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.line.color.rgb = fill
    write_text(shape, text, 8.5, bold=True, color=text_color)
    if target_slide is not None:
        shape.click_action.target_slide = target_slide
    elif hyperlink:
        shape.click_action.hyperlink.address = hyperlink
    return shape


def add_text_link(slide, text: str, left, top, width, height, target_slide, size: float = 13):
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    run = p.add_run()
    run.text = text
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    run.font.color.rgb = TEXT_DARK
    shape.click_action.target_slide = target_slide
    return shape


def excel_link(case: dict) -> str:
    # PowerPoint opens the workbook; supported Office builds also jump to the sheet/cell fragment.
    sheet = case["sheet"].replace("'", "''")
    return f"{SOURCE_XLSX.resolve().as_uri()}#'{sheet}'!A{case['row']}"


def add_toc(prs: Presentation, cases: list[dict], case_slides: list) -> object:
    toc_slide = prs.slides.add_slide(prs.slide_layouts[6])
    move_slide(prs, len(prs.slides) - 1, 1)

    bg = toc_slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    bg.fill.solid()
    bg.fill.fore_color.rgb = RGBColor(245, 247, 250)
    bg.line.color.rgb = RGBColor(245, 247, 250)

    title = toc_slide.shapes.add_textbox(Inches(0.6), Inches(0.35), Inches(12.1), Inches(0.65))
    write_text(title, "สารบัญ UAT Manufacturing", 28, bold=True, color=NAV_BLUE)

    subtitle = toc_slide.shapes.add_textbox(Inches(0.75), Inches(1.05), Inches(11.8), Inches(0.4))
    write_text(subtitle, "คลิกหัวข้อเพื่อไปยัง slide ของ scenario นั้น และใช้ปุ่มนำทางด้านล่างในแต่ละหน้า", 13, color=TEXT_DARK)

    sections: OrderedDict[str, dict] = OrderedDict()
    for idx, case in enumerate(cases):
        key = case["sheet"]
        if key not in sections:
            sections[key] = {
                "title": case["section"] or case["sheet"],
                "first_slide": case_slides[idx],
                "count": 0,
            }
        sections[key]["count"] += 1

    y = Inches(1.65)
    for number, (sheet, info) in enumerate(sections.items(), start=1):
        label = f"{number}. {info['title']} ({sheet}) - {info['count']} cases"
        add_text_link(toc_slide, label, Inches(1.0), y, Inches(11.4), Inches(0.38), info["first_slide"], size=13)
        y += Inches(0.48)

    add_button(
        toc_slide,
        "เปิดไฟล์ Excel UAT_MANU",
        Inches(9.15),
        Inches(6.9),
        Inches(2.7),
        Inches(0.38),
        hyperlink=SOURCE_XLSX.resolve().as_uri(),
        fill=NAV_GOLD,
        text_color=TEXT_DARK,
    )
    add_button(
        toc_slide,
        "เริ่ม Case แรก",
        Inches(6.55),
        Inches(6.9),
        Inches(2.35),
        Inches(0.38),
        target_slide=case_slides[0],
        fill=NAV_BLUE,
    )
    return toc_slide


def add_navigation(prs: Presentation, toc_slide, cases: list[dict], case_slides: list) -> None:
    slide_w = prs.slide_width
    slide_h = prs.slide_height
    top = slide_h - Inches(0.45)
    button_h = Inches(0.28)
    gap = Inches(0.06)
    widths = [Inches(0.85), Inches(0.85), Inches(0.85), Inches(0.85)]
    total_w = sum(widths) + gap * 3
    left = slide_w - total_w - Inches(0.35)

    add_button(
        prs.slides[0],
        "สารบัญ",
        slide_w - Inches(1.3),
        slide_h - Inches(0.55),
        Inches(0.95),
        Inches(0.32),
        target_slide=toc_slide,
        fill=NAV_GOLD,
        text_color=TEXT_DARK,
    )

    for idx, slide in enumerate(case_slides):
        x = left
        add_button(slide, "สารบัญ", x, top, widths[0], button_h, target_slide=toc_slide, fill=NAV_BLUE)
        x += widths[0] + gap
        previous_target = case_slides[idx - 1] if idx > 0 else toc_slide
        add_button(slide, "ก่อนหน้า", x, top, widths[1], button_h, target_slide=previous_target, fill=NAV_BLUE)
        x += widths[1] + gap
        next_target = case_slides[idx + 1] if idx < len(case_slides) - 1 else toc_slide
        add_button(slide, "ถัดไป", x, top, widths[2], button_h, target_slide=next_target, fill=NAV_BLUE)
        x += widths[2] + gap
        add_button(slide, "Excel", x, top, widths[3], button_h, hyperlink=excel_link(cases[idx]), fill=NAV_GOLD, text_color=TEXT_DARK)


def main() -> None:
    if not SOURCE_XLSX.exists():
        raise SystemExit(f"Source Excel not found: {SOURCE_XLSX}")
    if not INPUT_PPTX.exists():
        raise SystemExit(f"PowerPoint not found: {INPUT_PPTX}")

    cases = collect_cases()
    prs = Presentation(INPUT_PPTX)
    case_slides = list(prs.slides)[1:]
    if len(case_slides) != len(cases):
        raise SystemExit(f"Slide/case mismatch: slides={len(case_slides)} cases={len(cases)}")

    if INPUT_PPTX == OUTPUT_PPTX and not BACKUP_PPTX.exists():
        BACKUP_PPTX.write_bytes(INPUT_PPTX.read_bytes())

    toc_slide = add_toc(prs, cases, case_slides)
    add_navigation(prs, toc_slide, cases, case_slides)

    prs.save(OUTPUT_PPTX)
    print(OUTPUT_PPTX)
    print(f"slides={len(prs.slides)} cases={len(cases)} linked_navigation=ok")


if __name__ == "__main__":
    main()
