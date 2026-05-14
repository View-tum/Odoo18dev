from __future__ import annotations

import re
from copy import deepcopy
from pathlib import Path

from openpyxl import load_workbook
from pptx import Presentation
from pptx.enum.text import MSO_AUTO_SIZE
from pptx.util import Pt


TEMPLATE = Path(r"C:\Users\tumsu\Downloads\GMP_Accounting_Scenario_8Inventory Adj.pptx")
SOURCE = Path(r"C:\Users\tumsu\Downloads\UAT_MANU.xlsx")
OUTPUT = Path(r"C:\Users\tumsu\Downloads\GMP_Manufacturing_UAT_from_UAT_MANU.pptx")

ACTUAL_SHEET_PREFIX = tuple(f"{i:02d}_" for i in range(1, 10))


def norm(value) -> str:
    text = re.sub(r"\s+", " ", str(value or "").strip())
    text = text.replace("GMP/MOPH/...", "GMP/MOPH/เลขที่ MO ตัวอย่าง")
    text = text.replace("GMP/MOPL/...", "GMP/MOPL/เลขที่ MO ตัวอย่าง")
    text = text.replace("…", "")
    text = text.replace("...", "")
    return text


def shorten(text: str, limit: int) -> str:
    # Keep complete wording. The deck is generated for walkthrough/training,
    # so hiding text with ellipses is worse than using a smaller font.
    text = norm(text)
    return text


def split_steps(text: str, max_items: int = 5) -> list[str]:
    raw = str(text or "").replace("\r\n", "\n").replace("\r", "\n")
    parts: list[str] = []
    for line in raw.split("\n"):
        line = line.strip()
        if not line:
            continue
        line = re.sub(r"^\d+\)\s*", "", line)
        parts.append(line)
    if not parts and raw.strip():
        parts = [raw.strip()]
    return parts[:max_items]


def get_headers(ws) -> dict[str, int]:
    headers = {}
    for col in range(1, ws.max_column + 1):
        value = norm(ws.cell(6, col).value)
        if value:
            headers[value] = col
    return headers


def optional_cell(ws, row: int, headers: dict[str, int], name: str):
    idx = headers.get(name)
    if not idx:
        return None
    return ws.cell(row, idx).value


def collect_cases() -> list[dict]:
    wb = load_workbook(SOURCE, data_only=False)
    cases: list[dict] = []
    for ws in wb.worksheets:
        if not ws.title.startswith(ACTUAL_SHEET_PREFIX):
            continue
        headers = get_headers(ws)
        if "Case ID" not in headers or "Scenario ทดสอบ" not in headers:
            continue

        section_title = norm(ws["B2"].value) or ws.title
        objective = norm(ws["B3"].value)
        scope = norm(ws["B4"].value)

        for row in range(7, ws.max_row + 1):
            case_id = norm(ws.cell(row, headers["Case ID"]).value)
            if not re.match(r"^MU\d{2}-\d{2}$", case_id):
                continue
            cases.append(
                {
                    "sheet": ws.title,
                    "row": row,
                    "section": section_title,
                    "objective": objective,
                    "scope": scope,
                    "case_id": case_id,
                    "backlog": norm(ws.cell(row, headers.get("Backlog IDs", 3)).value),
                    "event": norm(ws.cell(row, headers.get("ลำดับเหตุการณ์", 4)).value),
                    "scenario": norm(ws.cell(row, headers.get("Scenario ทดสอบ", 5)).value),
                    "role": norm(optional_cell(ws, row, headers, "บทบาท / หน่วยงาน")),
                    "menu": norm(optional_cell(ws, row, headers, "Menu Path in local UAT (English)")),
                    "precondition": norm(optional_cell(ws, row, headers, "เงื่อนไขก่อนทดสอบ")),
                    "test_data": norm(optional_cell(ws, row, headers, "ข้อมูลทดสอบ")),
                    "steps": optional_cell(ws, row, headers, "ขั้นตอนทดสอบแบบละเอียด"),
                    "expected": norm(optional_cell(ws, row, headers, "ผลลัพธ์ที่คาดหวัง")),
                }
            )
    return cases


def duplicate_slide(prs: Presentation, source_slide):
    blank_layout = prs.slide_layouts[6]
    new_slide = prs.slides.add_slide(blank_layout)
    for shape in source_slide.shapes:
        new_slide.shapes._spTree.insert_element_before(
            deepcopy(shape._element),
            "p:extLst",
        )
    return new_slide


def delete_slide(prs: Presentation, index: int) -> None:
    slide_id_list = prs.slides._sldIdLst
    slide_id = slide_id_list[index]
    rel_id = slide_id.rId
    prs.part.drop_rel(rel_id)
    slide_id_list.remove(slide_id)


def set_text(shape, text: str, size: float = 12, bold: bool | None = None, italic: bool | None = None):
    tf = shape.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    p = tf.paragraphs[0]
    run = p.add_run()
    run.text = text
    run.font.name = "Tahoma"
    run.font.size = Pt(size)
    if bold is not None:
        run.font.bold = bold
    if italic is not None:
        run.font.italic = italic


def module_from_menu(menu: str) -> str:
    first = norm(menu).split(">")[0].strip()
    return f"Module {first}" if first else "Module Manufacturing"


def build_slide(slide, case: dict) -> None:
    shapes = list(slide.shapes)
    title = f"การผลิต: {shorten(case['scenario'], 45)}\n{case['case_id']}"
    set_text(shapes[0], title, size=18)

    menu = case["menu"] or "Manufacturing > Operations"
    # The template shape already has list numbering, so do not prefix "1)" here.
    set_text(shapes[1], shorten(menu, 95), size=11)
    set_text(shapes[9], module_from_menu(menu), size=12)

    middle_title = "ขั้นตอนทดสอบ"
    if case["event"]:
        middle_title = shorten(case["event"], 28)
    set_text(shapes[2], middle_title, size=11, italic=True)

    detail_lines = []
    if case["role"]:
        detail_lines.append(f"- ผู้ทดสอบ: {case['role']}")
    if case["precondition"]:
        detail_lines.append(f"- เงื่อนไข: {shorten(case['precondition'], 90)}")
    if case["test_data"]:
        detail_lines.append(f"- ข้อมูล: {shorten(case['test_data'], 90)}")
    if not detail_lines and case["objective"]:
        detail_lines.append(f"- วัตถุประสงค์: {shorten(case['objective'], 95)}")
    set_text(shapes[3], "\n".join(detail_lines[:4]), size=9.5)

    set_text(shapes[5], "ตรวจผลลัพธ์", size=12, bold=True, italic=True)

    steps = split_steps(case["steps"], max_items=5)
    step_text = "\n".join(f"- {shorten(item, 94)}" for item in steps)
    if case["expected"]:
        step_text += f"\n- ผลที่คาดหวัง: {shorten(case['expected'], 120)}"
    if case["backlog"]:
        step_text += f"\n- Backlog: {case['backlog']}"
    set_text(shapes[7], step_text.strip(), size=8.7)


def main() -> None:
    cases = collect_cases()
    if not cases:
        raise SystemExit("No UAT_MANU cases found")

    prs = Presentation(TEMPLATE)

    # Keep cover and one process slide; remove the second template process slide.
    while len(prs.slides) > 2:
        delete_slide(prs, 2)

    cover = prs.slides[0]
    set_text(cover.shapes[0], "UAT Manufacturing\nการผลิต", size=42, bold=True)

    base_process = prs.slides[1]
    for idx, case in enumerate(cases):
        slide = base_process if idx == 0 else duplicate_slide(prs, base_process)
        build_slide(slide, case)

    OUTPUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT)
    print(OUTPUT)
    print(f"slides={len(prs.slides)} cases={len(cases)}")


if __name__ == "__main__":
    main()
